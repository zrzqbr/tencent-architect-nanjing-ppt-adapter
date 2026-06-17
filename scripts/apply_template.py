#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
apply_template.py — 腾讯云架构师南京城市沙龙 PPT 模板适配器 v8（统一入口）

★ 自动识别输入格式，根据文件扩展名路由到对应管线：
  - .pptx         → PPTX 迁移适配（保留内容，替换品牌四件套）
  - .html/.htm    → HTML 转 PPT（解析结构，从零生成品牌合规 PPT）
  - .md/.markdown → Markdown 转 PPT（先转 HTML 再走 HTML 管线）

=== PPTX 迁移模式 ===
六项品牌适配：
  1) 替换背景图（按页面类型：封面/扉页/结尾页用 bg-cover/section/end/content 对应南京模板背景）
     - 插入前清除全屏遮罩（仅删除 z-order 底层的全屏图片/色块/渐变/图案，避免误删内容）
     - 清除 <p:bg> 层面的纯色/渐变/图案/图片填充
  2) 替换字体（中英文统一 TencentSans W3 正文 / TencentSans W7 标题，>=24pt 视为标题）
     - 递归处理组合形状（GroupShape）和表格单元格
  3) 添加 Logo（去重插入，按页面类型+输入PPT尺寸等比缩放定位）
  4) 配色合规：
     - 禁用色黑名单 → CIEDE2000 色差匹配 → 强制映射到品牌安全色
     - 主题色覆写（theme XML clrScheme 层面）
  5) 文字色对比度修复（P0-A）：
     - 背景替换后自动检测低对比度文字（WCAG 2.1 对比度 < 3:1）
     - 浅色背景+浅色文字 → 自动翻转为深蓝 #08194B
  6) 元素溢出检测（P1-B）：
     - 检测内容是否溢出页面边界或侵入南京页眉安全区（仅报告不修改）

=== HTML/Markdown 模式 ===
  - 解析 HTML 结构（H1-H6/p/ul/ol/table/pre+code/img/blockquote）
  - 智能分页 + 内容映射 → 直接生成品牌合规 PPT
  - 无需二次适配

使用：
  # PPTX 迁移
  python scripts/apply_template.py --input guest.pptx --output out.pptx [--mode full] [--dry-run]

  # HTML 转 PPT
  python scripts/apply_template.py --input article.html --output out.pptx [--title "标题"] [--author "作者"]

  # Markdown 转 PPT
  python scripts/apply_template.py --input notes.md --output out.pptx [--title "标题"]

依赖：
  pip install python-pptx Pillow beautifulsoup4 lxml
  （可选）pip install markdown  # Markdown 支持
"""

import argparse
import math
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
except ImportError:
    print("[ERROR] 缺少依赖。请运行: pip install python-pptx Pillow", file=sys.stderr)
    sys.exit(1)

# ---------- 配置 ----------
SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent
ASSETS = SKILL_DIR / "assets"

# 字体规范
FONT_TITLE_EA = "TencentSans W7"
FONT_BODY_EA = "TencentSans W3"
FONT_TITLE_LATIN = "TencentSans W7"
FONT_BODY_LATIN = "TencentSans W3"

CODE_FONT_KEYWORDS = ("Mono", "Consolas", "Menlo", "Courier", "Code")

# 标题字号阈值（单位：pt）
TITLE_SIZE_THRESHOLD_PT = 24    # >= 24pt 视为标题（W7 粗体）

# [P1-8] 扉页大标题字号阈值（单位：pt），用于 classify_page
SECTION_FONT_THRESHOLD_PT = 36  # >= 36pt 才视为扉页大标题（原硬编码 40pt 过严）

# [P0-1] 全屏遮罩检测阈值（遮盖幻灯片面积比例 ≥ 此值即视为全屏遮罩）
FULLSCREEN_OVERLAY_RATIO = 0.90


# ---------- 南京模板安全色板（限定 7 色）----------
# 用户指定：#3272DC、#08194B、#FFFFFF、#00C8D8、#01A4FF、#44474F、#8B8C8C
# 格式：(R, G, B, 描述)
BRAND_PALETTE = [
    (0x32, 0x72, 0xDC, "南京主蓝"),
    (0x08, 0x19, 0x4B, "南京深蓝"),
    (0xFF, 0xFF, 0xFF, "纯白"),
    (0x00, 0xC8, 0xD8, "青蓝强调"),
    (0x01, 0xA4, 0xFF, "亮蓝强调"),
    (0x44, 0x47, 0x4F, "深蓝辅助"),
    (0x8B, 0x8C, 0x8C, "中灰辅助"),
]

# 内容色块优先映射色。保留函数名以兼容旧调用，但不再使用旧暖色系。
WARM_BRAND_COLORS = [
    (0x32, 0x72, 0xDC, "南京主蓝"),
    (0x00, 0xC8, 0xD8, "青蓝强调"),
    (0x01, 0xA4, 0xFF, "亮蓝强调"),
    (0x08, 0x19, 0x4B, "南京深蓝"),
    (0x44, 0x47, 0x4F, "深蓝辅助"),
    (0x8B, 0x8C, 0x8C, "中灰辅助"),
]

# ---------- 南京模板禁用色黑名单 ----------
# 注意：#08194B 是南京核心深蓝，#3272DC/#00C8D8/#01A4FF 是南京主视觉蓝色，均不得禁用。
FORBIDDEN_COLOR_RULES = [
    # 旧红/暖色系在南京模板中禁用，避免残留旧模板风格
    {"name": "旧模板品牌红", "center": (0xD8, 0x0C, 0x01), "radius": 18},
    {"name": "纯红",     "center": (0xFF, 0x00, 0x00), "radius": 16},
    {"name": "暖杏",     "center": (0xD4, 0xA5, 0x74), "radius": 16},
    {"name": "米白",     "center": (0xFA, 0xF6, 0xEE), "radius": 10},
    {"name": "暖黄",     "center": (0xFA, 0xD1, 0x6A), "radius": 12},
    {"name": "提示金",   "center": (0xF1, 0x9D, 0x19), "radius": 12},
    # 与南京蓝色科技风冲突的高饱和色
    {"name": "鲜绿",     "center": (0x00, 0xFF, 0x00), "radius": 30},
    {"name": "草绿",     "center": (0x1D, 0x9A, 0x78), "radius": 18},
    {"name": "翠绿",     "center": (0x2E, 0xCC, 0x71), "radius": 18},
    {"name": "玫红",     "center": (0xFF, 0x14, 0x93), "radius": 22},
    {"name": "粉色",     "center": (0xFF, 0x69, 0xB4), "radius": 18},
    {"name": "紫色",     "center": (0x7C, 0x60, 0xC6), "radius": 18},
    {"name": "蓝紫",     "center": (0x7B, 0x61, 0xFF), "radius": 14},
    {"name": "荧光黄",   "center": (0xFF, 0xFF, 0x00), "radius": 15},
    {"name": "荧光橙",   "center": (0xFF, 0x7F, 0x00), "radius": 12},
]

# ---------- 色彩计算工具 ----------

def _rgb_to_linear(c: int) -> float:
    """sRGB 通道值 [0,255] → 线性光强度"""
    c_f = c / 255.0
    if c_f <= 0.04045:
        return c_f / 12.92
    return ((c_f + 0.055) / 1.055) ** 2.4


def _linear_to_xyz(r: int, g: int, b: int):
    """sRGB → CIE XYZ (D65)"""
    rl = _rgb_to_linear(r)
    gl = _rgb_to_linear(g)
    bl = _rgb_to_linear(b)
    x = rl * 0.4124564 + gl * 0.3575761 + bl * 0.1804375
    y = rl * 0.2126729 + gl * 0.7151522 + bl * 0.0721750
    z = rl * 0.0193339 + gl * 0.1191920 + bl * 0.9503041
    return x, y, z


def _xyz_to_lab(x: float, y: float, z: float):
    """CIE XYZ → CIE L*a*b* (D65 白点)"""
    xn, yn, zn = 0.95047, 1.00000, 1.08883

    def f(t):
        if t > 0.008856:
            return t ** (1/3)
        return 7.787 * t + 16/116

    l = 116 * f(y / yn) - 16
    a = 500 * (f(x / xn) - f(y / yn))
    b = 200 * (f(y / yn) - f(z / zn))
    return l, a, b


def rgb_to_lab(r: int, g: int, b: int):
    """RGB → CIE L*a*b*"""
    return _xyz_to_lab(*_linear_to_xyz(r, g, b))


def delta_e(r1: int, g1: int, b1: int, r2: int, g2: int, b2: int) -> float:
    """CIEDE2000 色差（P1-4 修复：替换原 CIE76 简化实现）。

    CIEDE2000 在蓝紫色区域的感知一致性显著优于 CIE76，
    避免将偏蓝的品牌色（如非南京色板蓝色）误判为禁用色。
    参考：Sharma et al. (2005), "The CIEDE2000 Color-Difference Formula"
    """
    L1, a1, b1_ = rgb_to_lab(r1, g1, b1)
    L2, a2, b2_ = rgb_to_lab(r2, g2, b2)

    # 步骤 1：计算 C'ab 和 h'ab
    C1 = math.sqrt(a1**2 + b1_**2)
    C2 = math.sqrt(a2**2 + b2_**2)
    C_avg = (C1 + C2) / 2.0

    # a' 调整系数 G（针对蓝紫区色相旋转）
    C_avg7 = C_avg**7
    G = 0.5 * (1 - math.sqrt(C_avg7 / (C_avg7 + 25**7)))
    a1p = a1 * (1 + G)
    a2p = a2 * (1 + G)

    C1p = math.sqrt(a1p**2 + b1_**2)
    C2p = math.sqrt(a2p**2 + b2_**2)

    def _atan2_deg(b_val, a_val):
        h = math.degrees(math.atan2(b_val, a_val))
        return h + 360 if h < 0 else h

    h1p = _atan2_deg(b1_, a1p) if C1p > 1e-9 else 0.0
    h2p = _atan2_deg(b2_, a2p) if C2p > 1e-9 else 0.0

    # 步骤 2：ΔL' ΔC' Δh'
    dLp = L2 - L1
    dCp = C2p - C1p

    if C1p * C2p < 1e-9:
        dhp = 0.0
    elif abs(h2p - h1p) <= 180:
        dhp = h2p - h1p
    elif h2p - h1p > 180:
        dhp = h2p - h1p - 360
    else:
        dhp = h2p - h1p + 360

    dHp = 2 * math.sqrt(C1p * C2p) * math.sin(math.radians(dhp / 2))

    # 步骤 3：均值
    Lp_avg = (L1 + L2) / 2.0
    Cp_avg = (C1p + C2p) / 2.0

    if C1p * C2p < 1e-9:
        hp_avg = h1p + h2p
    elif abs(h1p - h2p) <= 180:
        hp_avg = (h1p + h2p) / 2.0
    elif h1p + h2p < 360:
        hp_avg = (h1p + h2p + 360) / 2.0
    else:
        hp_avg = (h1p + h2p - 360) / 2.0

    T = (1
         - 0.17 * math.cos(math.radians(hp_avg - 30))
         + 0.24 * math.cos(math.radians(2 * hp_avg))
         + 0.32 * math.cos(math.radians(3 * hp_avg + 6))
         - 0.20 * math.cos(math.radians(4 * hp_avg - 63)))

    SL = 1 + 0.015 * (Lp_avg - 50)**2 / math.sqrt(20 + (Lp_avg - 50)**2)
    SC = 1 + 0.045 * Cp_avg
    SH = 1 + 0.015 * Cp_avg * T

    Cp_avg7 = Cp_avg**7
    RC = 2 * math.sqrt(Cp_avg7 / (Cp_avg7 + 25**7))
    d_theta = 30 * math.exp(-((hp_avg - 275) / 25)**2)
    RT = -math.sin(math.radians(2 * d_theta)) * RC

    kL = kC = kH = 1.0
    return math.sqrt(
        (dLp / (kL * SL))**2
        + (dCp / (kC * SC))**2
        + (dHp / (kH * SH))**2
        + RT * (dCp / (kC * SC)) * (dHp / (kH * SH))
    )


def _is_forbidden_color(r: int, g: int, b: int) -> bool:
    """检查颜色是否在禁用色黑名单内（蓝紫/青/鲜绿等）。"""
    for rule in FORBIDDEN_COLOR_RULES:
        cr, cg, cb = rule["center"]
        de = delta_e(r, g, b, cr, cg, cb)
        if de <= rule["radius"]:
            return True
    return False


def find_nearest_brand_color(r: int, g: int, b: int,
                              max_delta: float = 30.0,
                              skip_white_black: bool = False,
                              force_replace: bool = False):
    """在品牌色板中找最近的颜色。

    逻辑：
    1. 若颜色命中禁用色黑名单 → 强制替换（max_delta 放宽到 80，找最近品牌色）
    2. 否则若 skip_white_black=True 且为中性色 → 跳过
    3. 否则按 max_delta 正常匹配

    返回 (R, G, B, 描述, ΔE) 或 None（跳过）。
    force_replace=True 时跳过中性色检查（用于图表配色）。
    """
    is_forbidden = _is_forbidden_color(r, g, b)

    # 中性色判定（饱和度 < 12%）
    max_ch = max(r, g, b)
    min_ch = min(r, g, b)
    saturation = (max_ch - min_ch) / max(max_ch, 1)
    is_neutral = saturation < 0.12

    # 中性色且不是禁用色 → 跳过（避免把 #262626 等正文近黑色误改）
    # 但如果是 force_replace 模式（图表），还是替换
    if skip_white_black and is_neutral and not is_forbidden and not force_replace:
        return None

    # 禁用色：用更宽松阈值
    effective_max_delta = 80.0 if is_forbidden else max_delta

    best = None
    best_de = float('inf')
    for (br, bg, bb, desc) in BRAND_PALETTE:
        de = delta_e(r, g, b, br, bg, bb)
        if de < best_de:
            best_de = de
            best = (br, bg, bb, desc, de)

    if best and best_de <= effective_max_delta:
        return best
    return None


def find_nearest_warm_color(r: int, g: int, b: int):
    """兼容旧调用：在南京内容色块优先色中找最近的颜色。

    当前已不再使用旧暖色体系；该函数只从南京限定色板中选择。
    返回 (R, G, B, 描述, ΔE) 或 None。
    """
    best = None
    best_de = float('inf')
    for (wr, wg, wb, desc) in WARM_BRAND_COLORS:
        de = delta_e(r, g, b, wr, wg, wb)
        if de < best_de:
            best_de = de
            best = (wr, wg, wb, desc, de)
    return best


def _rgb_unpack(rgb) -> tuple:
    """解包 pptx RGBColor 对象（实际上是支持索引的 int 子类）→ (r, g, b)。"""
    return (rgb[0], rgb[1], rgb[2])


def _get_rgb_from_solid_fill(fill) -> tuple:
    """从 pptx FillFormat 提取 RGB，返回 (r, g, b) 或 None。

    只处理 SOLID 填充且颜色为 RGB 类型的情况，其余（主题色、继承色等）跳过。
    """
    try:
        ft = fill.type
        if ft is None:
            return None
        if ft.name != 'SOLID':
            return None
        fc = fill.fore_color
        if fc.type is None:
            return None
        if fc.type.name != 'RGB':
            return None
        return _rgb_unpack(fc.rgb)
    except Exception:
        pass
    return None


def _set_solid_fill_rgb(fill, r: int, g: int, b: int):
    """把形状的 solid fill 改为指定 RGB。"""
    try:
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
    except Exception:
        pass


def _get_run_text_color_xml(run) -> tuple:
    """通过 XML 直接读取 text run 的 sRGB 颜色，返回 (r, g, b) 或 None。

    pptx 的 font.color.type 在某些场合会误报，直接查 XML 更可靠。
    同时支持解析 schemeClr 主题引用（需传入 theme_cache）。
    """
    try:
        rPr = run._r.find(qn('a:rPr'))
        if rPr is None:
            return None
        solidFill = rPr.find('.//' + qn('a:solidFill'))
        if solidFill is None:
            return None

        # 优先读取显式 sRGB
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            hex_val = srgbClr.get('val', '')
            if len(hex_val) == 6:
                r = int(hex_val[0:2], 16)
                g = int(hex_val[2:4], 16)
                b = int(hex_val[4:6], 16)
                return (r, g, b)

        # 尝试解析主题引用 schemeClr
        schemeClr = solidFill.find(qn('a:schemeClr'))
        if schemeClr is not None:
            slot = schemeClr.get('val', '')
            if slot and _theme_color_cache:
                # 应用亮度/饱和度调节（shade/tint/lumMod/lumOff）
                resolved = _resolve_theme_slot_with_modifiers(schemeClr, slot)
                if resolved:
                    return resolved
    except Exception:
        pass
    return None


# ---------- 主题色引用解析缓存 ----------
_theme_color_cache: dict = {}  # slot_name → (r, g, b)


def _build_theme_color_cache(pres):
    """解析 PPT 中所有主题的 clrScheme，构建 slot → RGB 映射。

    处理多个 slide master 和多个主题 part，合并后存入模块全局缓存。
    """
    _theme_color_cache.clear()
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    SLOT_NAMES = ["dk1", "lt1", "dk2", "lt2",
                  "accent1", "accent2", "accent3", "accent4",
                  "accent5", "accent6", "hlink", "folHlink"]
    processed = set()

    for master in pres.slide_masters:
        try:
            for rel in master.part.rels.values():
                if 'theme' not in str(rel.reltype).lower():
                    continue
                theme_part = rel.target_part
                if id(theme_part) in processed:
                    continue
                processed.add(id(theme_part))

                from lxml import etree as ET
                root = ET.fromstring(theme_part.blob)
                ns_map = {'a': ns}
                clrScheme = root.find('.//a:clrScheme', ns_map)
                if clrScheme is None:
                    continue
                for child in clrScheme:
                    slot_local = ET.QName(child).localname
                    if slot_local not in SLOT_NAMES:
                        continue
                    for color_el in child:
                        if ET.QName(color_el).localname == 'srgbClr':
                            val = color_el.get('val', '')
                            if len(val) == 6:
                                r = int(val[0:2], 16)
                                g = int(val[2:4], 16)
                                b = int(val[4:6], 16)
                                _theme_color_cache[slot_local] = (r, g, b)
                            break
        except Exception:
            continue


def _resolve_theme_slot_with_modifiers(scheme_clr_el, slot: str) -> tuple:
    """解析 theme schemeClr + 可选 shade/tint/lumMod/lumOff，返回 (r, g, b)。

    PowerPoint OOXML 的 shade/tint 通过线性插值实现：
    - shade: 向黑色混合 (factor × black)
    - tint: 向白色混合 (factor × white)

    OOXML 别名映射：bg1→lt1, bg2→lt2, tx1→dk1, tx2→dk2
    """
    # OOXML schemeClr 别名映射
    SLOT_ALIASES = {
        "bg1": "lt1", "bg2": "lt2",
        "tx1": "dk1", "tx2": "dk2",
    }
    resolved_slot = SLOT_ALIASES.get(slot, slot)
    base = _theme_color_cache.get(resolved_slot)
    if base is None:
        return None
    r, g, b = base

    # lumMod（亮度缩放，0-100000 表示 0%-1000%）
    lumMod_el = scheme_clr_el.find(qn('a:lumMod'))
    if lumMod_el is not None:
        try:
            factor = int(lumMod_el.get('val', '100000')) / 100000.0
            r = min(255, int(r * factor))
            g = min(255, int(g * factor))
            b = min(255, int(b * factor))
        except (ValueError, TypeError):
            pass

    # lumOff（亮度偏移，0-100000 表示 0%-1000%）
    lumOff_el = scheme_clr_el.find(qn('a:lumOff'))
    if lumOff_el is not None:
        try:
            offset = int(lumOff_el.get('val', '0')) / 100000.0
            r = min(255, int(r + offset * 255))
            g = min(255, int(g + offset * 255))
            b = min(255, int(b + offset * 255))
        except (ValueError, TypeError):
            pass

    # shade（变暗，factor * black）
    shade_el = scheme_clr_el.find(qn('a:shade'))
    if shade_el is not None:
        try:
            factor = int(shade_el.get('val', '100000')) / 100000.0
            r = int(r * factor)
            g = int(g * factor)
            b = int(b * factor)
        except (ValueError, TypeError):
            pass

    # tint（变亮，factor * white + (1-factor) * base）
    tint_el = scheme_clr_el.find(qn('a:tint'))
    if tint_el is not None:
        try:
            factor = int(tint_el.get('val', '0')) / 100000.0
            r = int(r + (255 - r) * factor)
            g = int(g + (255 - g) * factor)
            b = int(b + (255 - b) * factor)
        except (ValueError, TypeError):
            pass

    return (min(255, max(0, r)), min(255, max(0, g)), min(255, max(0, b)))


def _set_run_text_color_xml(run, r: int, g: int, b: int):
    """通过 XML 直接修改 text run 的 sRGB 颜色。"""
    try:
        from lxml import etree
        rPr = run._r.get_or_add_rPr()
        # 先删除旧的颜色节点
        for tag in (qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'), qn('a:pattFill')):
            old = rPr.find(tag)
            if old is not None:
                rPr.remove(old)
        # 写入新的 solidFill/srgbClr
        solid = etree.SubElement(rPr, qn('a:solidFill'))
        srgb = etree.SubElement(solid, qn('a:srgbClr'))
        srgb.set('val', f'{r:02X}{g:02X}{b:02X}')
    except Exception:
        pass


def _fix_run_scheme_white(run, dry_run: bool = False) -> bool:
    """检查并修复 run 中 rPr/solidFill/schemeClr 引用白色槽位（bg1/lt1/bg2/lt2）的情况。

    如果检测到白色 scheme 引用，将其替换为 srgbClr #08194B。
    返回 True 表示进行了修复，False 表示无需修复。
    """
    try:
        from lxml import etree as _etree
        rPr = run._r.find(qn('a:rPr'))
        if rPr is None:
            return False
        solidFill = rPr.find(qn('a:solidFill'))
        if solidFill is None:
            return False
        schemeClr = solidFill.find(qn('a:schemeClr'))
        if schemeClr is None:
            return False
        slot = schemeClr.get('val', '')
        WHITE_SLOTS = ('bg1', 'lt1', 'bg2', 'lt2')
        if slot not in WHITE_SLOTS:
            return False
        if not dry_run:
            # 删除 schemeClr 及其子元素（lumMod/lumOff等），替换为 srgbClr
            solidFill.remove(schemeClr)
            new_srgb = _etree.SubElement(solidFill, qn('a:srgbClr'))
            new_srgb.set('val', '08194B')
        return True
    except Exception:
        return False


def _get_line_rgb_xml(shp_element) -> tuple:
    """通过 XML 直接读取形状边框线条颜色，返回 (r, g, b) 或 None。"""
    try:
        for ln_el in shp_element.iter(qn('a:ln')):
            for sf in ln_el.iter(qn('a:solidFill')):
                for srgb in sf.iter(qn('a:srgbClr')):
                    val = srgb.get('val', '')
                    if len(val) == 6:
                        return (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
    except Exception:
        pass
    return None


def _set_line_rgb_xml(shp_element, r: int, g: int, b: int):
    """通过 XML 直接修改形状边框线条颜色。"""
    try:
        for ln_el in shp_element.iter(qn('a:ln')):
            for sf in list(ln_el.iter(qn('a:solidFill'))):
                for srgb in list(sf.iter(qn('a:srgbClr'))):
                    srgb.set('val', f'{r:02X}{g:02X}{b:02X}')
                    return  # 只改第一个
    except Exception:
        pass


def _get_line_width_xml(shp_element) -> float:
    """读取形状边框宽度，返回 pt 值。0 表示无边框/隐形。"""
    try:
        for ln_el in shp_element.iter(qn('a:ln')):
            w_attr = ln_el.get('w')
            if w_attr is not None:
                # EMU 单位，12700 EMU = 1pt
                return int(w_attr) / 12700.0
            # 没有 w 属性 → 使用默认宽度（PowerPoint 默认 0.75pt）
            # 只要有 solidFill 就说明有边框
            for sf in ln_el.iter(qn('a:solidFill')):
                return 0.75  # 默认
    except Exception:
        pass
    return 0.0


def _remove_line_xml(shp_element):
    """通过 XML 移除形状边框（设置 noFill）。"""
    try:
        from lxml import etree
        for ln_el in shp_element.iter(qn('a:ln')):
            # 清空所有填充子节点，改为 noFill
            for child in list(ln_el):
                if child.tag in (qn('a:solidFill'), qn('a:gradFill'), qn('a:pattFill')):
                    ln_el.remove(child)
            # 添加 noFill
            if ln_el.find(qn('a:noFill')) is None:
                etree.SubElement(ln_el, qn('a:noFill'))
            return
    except Exception:
        pass


def replace_colors_in_slide(slide, max_delta: float = 30.0,
                             skip_white_black: bool = False,
                             dry_run: bool = False) -> list:
    """扫描并替换当前幻灯片中不合规的颜色。

    改进点（v5）：
    - 使用 CIEDE2000 色差（P1-4），蓝紫边界判断更准确
    - 禁用色（蓝紫/青/鲜绿等）使用宽松阈值强制替换
    - 边框颜色：只处理实际可见边框，不激活宽度=0的隐形边框
    - 当边框颜色是禁用色时，将边框颜色替换为品牌色（保留边框结构）
    - 当边框颜色是黑色/白色等中性品牌色时，保持不变

    返回替换记录列表：[{"element": str, "from": hex, "to": hex, "delta_e": float, "brand_name": str}]
    """
    changes = []

    def log_change(element: str, orig_rgb: tuple, new_rgb: tuple, de: float, desc: str):
        changes.append({
            "element": element,
            "from": "#{:02X}{:02X}{:02X}".format(*orig_rgb),
            "to": "#{:02X}{:02X}{:02X}".format(*new_rgb),
            "delta_e": round(de, 1),
            "brand_name": desc,
        })

    for shp in slide.shapes:
        shape_name = getattr(shp, 'name', str(shp.shape_id))

        # 1) 形状填充色（SOLID RGB 填充）— 暖色优先：内容色块避免全部映射到商务蓝
        try:
            fill = shp.fill
            orig = _get_rgb_from_solid_fill(fill)
            if orig:
                nearest = find_nearest_brand_color(*orig, max_delta=max_delta,
                                                   skip_white_black=skip_white_black)
                if nearest:
                    nr, ng, nb, desc, de = nearest
                    if (nr, ng, nb) != orig:
                        log_change(f"形状[{shape_name}].填充", orig, (nr, ng, nb), de, desc)
                        if not dry_run:
                            _set_solid_fill_rgb(fill, nr, ng, nb)
        except Exception:
            pass

        # 2) 线条/边框色（只处理实际可见边框）
        try:
            sp_el = shp._element
            line_rgb = _get_line_rgb_xml(sp_el)
            if line_rgb:
                line_width = _get_line_width_xml(sp_el)
                is_forbidden = _is_forbidden_color(*line_rgb)

                if is_forbidden and line_width > 0:
                    # 禁用色且有实际宽度 → 替换为最近品牌色
                    nearest = find_nearest_brand_color(*line_rgb, max_delta=80.0,
                                                       skip_white_black=False)
                    if nearest:
                        nr, ng, nb, desc, de = nearest
                        if (nr, ng, nb) != line_rgb:
                            log_change(f"形状[{shape_name}].边框", line_rgb, (nr, ng, nb), de, desc)
                            if not dry_run:
                                _set_line_rgb_xml(sp_el, nr, ng, nb)
                elif is_forbidden and line_width == 0:
                    # 禁用色但宽度=0（隐形边框） → 直接移除边框定义，避免意外激活
                    log_change(f"形状[{shape_name}].边框(隐形→移除)", line_rgb, (0, 0, 0), 0, "移除禁用色隐形边框")
                    if not dry_run:
                        _remove_line_xml(sp_el)
                elif not is_forbidden and line_width > 0:
                    # 非禁用色但有实际边框 → 按标准 ΔE 阈值处理
                    nearest = find_nearest_brand_color(*line_rgb, max_delta=max_delta,
                                                       skip_white_black=skip_white_black)
                    if nearest:
                        nr, ng, nb, desc, de = nearest
                        if (nr, ng, nb) != line_rgb:
                            log_change(f"形状[{shape_name}].边框", line_rgb, (nr, ng, nb), de, desc)
                            if not dry_run:
                                _set_line_rgb_xml(sp_el, nr, ng, nb)
                # else: 非禁用色 + 宽度=0 → 完全跳过，不触碰
        except Exception:
            pass

        # 3) 文字颜色（通过 XML 直读，比 pptx API 更可靠）
        if shp.has_text_frame:
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        orig = _get_run_text_color_xml(run)
                        if orig:
                            nearest = find_nearest_brand_color(*orig, max_delta=max_delta,
                                                               skip_white_black=skip_white_black)
                            if nearest:
                                nr, ng, nb, desc, de = nearest
                                if (nr, ng, nb) != orig:
                                    log_change(f"形状[{shape_name}].文字", orig, (nr, ng, nb), de, desc)
                                    if not dry_run:
                                        _set_run_text_color_xml(run, nr, ng, nb)
                    except Exception:
                        pass

        # 4) 图表配色（图表系列不跳过中性色，全量替换）
        if shp.has_chart:
            try:
                chart = shp.chart
                for series in chart.series:
                    try:
                        fill = series.format.fill
                        orig = _get_rgb_from_solid_fill(fill)
                        if orig:
                            nearest = find_nearest_brand_color(*orig, max_delta=max_delta,
                                                               skip_white_black=False,
                                                               force_replace=True)
                            if nearest:
                                nr, ng, nb, desc, de = nearest
                                if (nr, ng, nb) != orig:
                                    log_change(f"图表[{shape_name}].系列填充", orig, (nr, ng, nb), de, desc)
                                    if not dry_run:
                                        _set_solid_fill_rgb(fill, nr, ng, nb)
                    except Exception:
                        pass
            except Exception:
                pass

        # 5) 表格单元格填充 — 暖色优先：表头/数据格避免全部映射到商务蓝
        if shp.has_table:
            try:
                table = shp.table
                for row in table.rows:
                    for cell in row.cells:
                        try:
                            fill = cell.fill
                            orig = _get_rgb_from_solid_fill(fill)
                            if orig:
                                nearest = find_nearest_brand_color(*orig, max_delta=max_delta,
                                                                   skip_white_black=skip_white_black)
                                if nearest:
                                    nr, ng, nb, desc, de = nearest
                                    if (nr, ng, nb) != orig:
                                        log_change(f"表格[{shape_name}].单元格填充", orig, (nr, ng, nb), de, desc)
                                        if not dry_run:
                                            _set_solid_fill_rgb(fill, nr, ng, nb)
                        except Exception:
                            pass
            except Exception:
                pass

        # 6) 组合形状递归处理（P1-8 改进：支持子形状扫描）
        try:
            if hasattr(shp, 'shapes'):  # GroupShapes
                for sub_shp in shp.shapes:
                    sub_name = getattr(sub_shp, 'name', str(sub_shp.shape_id))
                    # 递归处理子形状填充 — 暖色优先
                    try:
                        fill = sub_shp.fill
                        orig = _get_rgb_from_solid_fill(fill)
                        if orig:
                            nearest = find_nearest_brand_color(*orig, max_delta=max_delta,
                                                               skip_white_black=skip_white_black)
                            if nearest:
                                nr, ng, nb, desc, de = nearest
                                if (nr, ng, nb) != orig:
                                    log_change(f"组合[{shape_name}].子[{sub_name}].填充", orig, (nr, ng, nb), de, desc)
                                    if not dry_run:
                                        _set_solid_fill_rgb(fill, nr, ng, nb)
                    except Exception:
                        pass
                    # 递归处理子形状边框
                    try:
                        sp_el = sub_shp._element
                        line_rgb = _get_line_rgb_xml(sp_el)
                        if line_rgb:
                            is_forbidden = _is_forbidden_color(*line_rgb)
                            line_width = _get_line_width_xml(sp_el)
                            if is_forbidden:
                                if line_width > 0:
                                    nearest = find_nearest_brand_color(*line_rgb, max_delta=80.0,
                                                                       skip_white_black=False)
                                    if nearest:
                                        nr, ng, nb, desc, de = nearest
                                        if (nr, ng, nb) != line_rgb:
                                            log_change(f"组合[{shape_name}].子[{sub_name}].边框", line_rgb, (nr, ng, nb), de, desc)
                                            if not dry_run:
                                                _set_line_rgb_xml(sp_el, nr, ng, nb)
                                else:
                                    log_change(f"组合[{shape_name}].子[{sub_name}].边框(隐形→移除)", line_rgb, (0, 0, 0), 0, "移除禁用色隐形边框")
                                    if not dry_run:
                                        _remove_line_xml(sp_el)
                    except Exception:
                        pass
                    # 递归处理子形状文字
                    if sub_shp.has_text_frame:
                        for para in sub_shp.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    orig = _get_run_text_color_xml(run)
                                    if orig:
                                        nearest = find_nearest_brand_color(*orig, max_delta=max_delta,
                                                                           skip_white_black=skip_white_black)
                                        if nearest:
                                            nr, ng, nb, desc, de = nearest
                                            if (nr, ng, nb) != orig:
                                                log_change(f"组合[{shape_name}].子[{sub_name}].文字", orig, (nr, ng, nb), de, desc)
                                                if not dry_run:
                                                    _set_run_text_color_xml(run, nr, ng, nb)
                                except Exception:
                                    pass
        except Exception:
            pass

    return changes


# ---------- 主题色覆写（v4 新增）----------

def replace_theme_colors(pres, dry_run: bool = False) -> list:
    """覆写 PPT 主题色定义中的禁用色。

    PPT 中形状样式（sp:style/fillRef 等）引用 accent1/dk2 等主题槽时，
    渲染颜色由主题定义决定，srgbClr 扫描无法捕捉。
    本函数直接修改 theme XML 中的 clrScheme，将禁用色主题槽替换为最近品牌色。

    返回替换记录列表。
    """
    changes = []
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # 品牌色槽映射：主题槽名 → 期望的品牌色（若原来是禁用色则覆写）
    # 同时也会按 find_nearest_brand_color 智能选择
    SLOT_NAMES = ["dk1", "lt1", "dk2", "lt2",
                  "accent1", "accent2", "accent3", "accent4",
                  "accent5", "accent6", "hlink", "folHlink"]

    processed_theme_parts = set()

    for master in pres.slide_masters:
        try:
            for rel in master.part.rels.values():
                if 'theme' not in str(rel.reltype).lower():
                    continue
                theme_part = rel.target_part
                part_id = id(theme_part)
                if part_id in processed_theme_parts:
                    continue
                processed_theme_parts.add(part_id)

                from lxml import etree as ET
                xml_bytes = theme_part.blob
                root = ET.fromstring(xml_bytes)
                ns_map = {'a': ns}
                clrScheme = root.find('.//a:clrScheme', ns_map)
                if clrScheme is None:
                    continue

                modified = False
                for child in clrScheme:
                    slot_local = ET.QName(child).localname
                    if slot_local not in SLOT_NAMES:
                        continue
                    for color_el in list(child):
                        color_local = ET.QName(color_el).localname
                        if color_local != 'srgbClr':
                            continue
                        val = color_el.get('val', '')
                        if len(val) != 6:
                            continue
                        r = int(val[0:2], 16)
                        g = int(val[2:4], 16)
                        b = int(val[4:6], 16)

                        # 只替换禁用色
                        if not _is_forbidden_color(r, g, b):
                            continue

                        nearest = find_nearest_brand_color(r, g, b,
                                                           max_delta=80.0,
                                                           skip_white_black=False)
                        if nearest is None:
                            continue
                        nr, ng, nb, desc, de = nearest
                        if (nr, ng, nb) == (r, g, b):
                            continue

                        old_hex = f"#{val.upper()}"
                        new_hex = f"#{nr:02X}{ng:02X}{nb:02X}"
                        changes.append({
                            "element": f"主题[{slot_local}]",
                            "from": old_hex,
                            "to": new_hex,
                            "delta_e": round(de, 1),
                            "brand_name": desc,
                            "slide": 0,  # 主题级，非页面级
                        })

                        if not dry_run:
                            color_el.set('val', f"{nr:02X}{ng:02X}{nb:02X}")
                            modified = True

                if modified and not dry_run:
                    # 写回修改后的 XML 到 theme part
                    new_xml = ET.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                    theme_part._blob = new_xml

        except Exception as e:
            # [P1-6] 异常不再静默吞噬，由调用方写入 warnings
            raise

    return changes


# ---------- 资产获取 ----------

def get_assets(city: str = "nanjing"):
    """返回南京模板背景路径。南京模板不再自动插入旧横版 Logo。"""
    def resolve_bg(name: str) -> Path:
        generic = ASSETS / "backgrounds" / name
        if generic.exists():
            return generic
        city_specific = ASSETS / "backgrounds" / f"{city}-{name}"
        if city_specific.exists():
            return city_specific
        return generic

    bg_cover = resolve_bg("bg-cover.jpeg")
    bg_section = resolve_bg("bg-section.jpeg")
    bg_content = resolve_bg("bg-content.jpeg")
    bg_end = resolve_bg("bg-end.jpeg")
    bg_hero_cover = resolve_bg("bg-hero-cover.jpeg")

    for p in (bg_cover, bg_section, bg_content, bg_end):
        if not p.exists():
            raise FileNotFoundError(f"缺少南京模板资源：{p}")

    return {
        "bg_cover": str(bg_cover),
        "bg_section": str(bg_section),
        "bg_content": str(bg_content),
        "bg_end": str(bg_end),
        "bg_hero_cover": str(bg_hero_cover) if bg_hero_cover.exists() else str(bg_cover),
        "logo_main": None,
        "logo_corner": None,
        "logo_enabled": False,
        "city": city,
    }


# ---------- 页面类型识别（P1-8 改进）----------

def _collect_text_shapes(slide) -> tuple:
    """递归收集幻灯片中所有文字形状（包含组合内子形状）。

    返回 (text_shapes_count, big_title_count, has_chart, has_table)。
    big_title 阈值使用 SECTION_FONT_THRESHOLD_PT（可配置常量，默认 36pt）。
    """
    text_shapes = 0
    has_chart = False
    has_table = False
    big_title_count = 0

    def _scan_shape(shp):
        nonlocal text_shapes, has_chart, has_table, big_title_count

        if shp.has_text_frame and shp.text_frame.text.strip():
            text_shapes += 1
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    # [P1-8] 字号缺失时回退为 0，不用 18 假设
                    size_pt = run.font.size.pt if run.font.size else 0
                    if size_pt >= SECTION_FONT_THRESHOLD_PT:
                        big_title_count += 1
        if shp.has_chart:
            has_chart = True
        if shp.has_table:
            has_table = True

        # [P1-8] 递归进入组合形状
        if hasattr(shp, 'shapes'):
            for sub in shp.shapes:
                _scan_shape(sub)

    for shp in slide.shapes:
        _scan_shape(shp)

    return text_shapes, big_title_count, has_chart, has_table


def classify_page(slide, idx: int, total: int) -> str:
    """判断页面类型：cover / section / content / end。

    v7 改进（P1-A）：
    - 优先级 1：从 slide_layout.name 匹配（最可靠）
    - 优先级 2：首尾页默认规则
    - 优先级 3：内容启发式（文本数量+大标题字号+图表）
    """
    # 优先级 1：Layout 名称匹配
    try:
        layout_name = slide.slide_layout.name.lower()
        # 封面关键词
        cover_kw = ("封面", "cover", "title slide", "首页", "封底")
        if any(k in layout_name for k in cover_kw):
            # 最后一页的"封底"类 layout 归为 end
            if idx == total - 1 or "封底" in layout_name:
                return "end"
            return "cover"
        # 章节/扉页关键词
        section_kw = ("章节", "section", "扉页", "divider", "break")
        if any(k in layout_name for k in section_kw):
            return "section"
        # 结尾关键词
        end_kw = ("结尾", "ending", "结束", "thank", "end")
        if any(k in layout_name for k in end_kw):
            return "end"
    except Exception:
        pass

    # 优先级 2：首尾页默认
    if idx == 0:
        return "cover"
    if idx == total - 1:
        return "end"

    # 优先级 3：内容启发式
    text_shapes, big_title_count, has_chart, has_table = _collect_text_shapes(slide)

    if text_shapes <= 2 and big_title_count >= 1 and not (has_chart or has_table):
        return "section"
    return "content"


# ---------- P0-1：全屏遮罩清除 ----------

def remove_fullscreen_overlays(slide, slide_w: int, slide_h: int,
                                dry_run: bool = False) -> int:
    """检测并删除幻灯片中覆盖全页的不透明遮罩形状（矩形/图片）。

    判断标准：形状的宽度 ≥ slide_w × FULLSCREEN_OVERLAY_RATIO
              且高度 ≥ slide_h × FULLSCREEN_OVERLAY_RATIO
              且填充为 SOLID（不透明色块）或为图片（p:pic）

    被删除的形状会记录到日志中。返回删除数量。

    典型场景：原 PPT 用一张全屏图片或色块作为背景→运行适配后插入新背景图→
    两层叠加导致原始内容被遮住。本函数在插入新背景前调用，清除旧遮罩。
    """
    from lxml import etree

    min_w = slide_w * FULLSCREEN_OVERLAY_RATIO
    min_h = slide_h * FULLSCREEN_OVERLAY_RATIO

    to_remove = []

    spTree = slide.shapes._spTree
    shape_elements = [child for child in spTree.iterchildren()
                      if etree.QName(child).localname in ("sp", "pic", "grpSp", "graphicFrame", "cxnSp")]

    for shp in slide.shapes:
        try:
            # 检查尺寸（EMU）
            if shp.width < min_w or shp.height < min_h:
                continue
            # 位置也要在幻灯片范围内（允许少量偏移）
            margin = slide_w * 0.05
            if shp.left > margin or shp.top > margin:
                continue

            sp_el = shp._element
            tag_local = etree.QName(sp_el).localname

            # 图片形状（p:pic）：仅删除 z-order 最底层的全屏图片（视为旧背景）
            # 避免误删用户有意保留的全屏内容图片（如截图、海报）
            if tag_local == 'pic':
                # 只有位于 spTree 前 2 个形状位置的全屏图片才视为背景遮罩
                idx_in_tree = shape_elements.index(sp_el) if sp_el in shape_elements else 999
                if idx_in_tree <= 1:
                    to_remove.append((shp, sp_el, "图片遮罩（z-order底层）"))
                continue

            # 形状（p:sp）：检查是否为不透明填充（实心/渐变/图案）
            if tag_local == 'sp':
                try:
                    fill = shp.fill
                    ft = fill.type
                    if ft is not None and ft.name == 'SOLID':
                        to_remove.append((shp, sp_el, "实心色块遮罩"))
                    elif ft is not None and ft.name == 'GRADIENT':
                        to_remove.append((shp, sp_el, "渐变填充遮罩"))
                    elif ft is not None and ft.name == 'PATTERNED':
                        to_remove.append((shp, sp_el, "图案填充遮罩"))
                    elif ft is not None and ft.name == 'BACKGROUND':
                        to_remove.append((shp, sp_el, "背景填充遮罩"))
                except Exception:
                    pass
        except Exception:
            pass

    removed = 0
    for shp, sp_el, reason in to_remove:
        try:
            if not dry_run:
                sp_el.getparent().remove(sp_el)
            removed += 1
        except Exception:
            pass

    return removed


# ---------- 背景替换 ----------

def set_slide_background(slide, image_path: str):
    """用图片填充幻灯片背景，发送到 z-order 最底层。

    v6 修复：先清除上游 Skill 可能写入的 <p:bg> 纯色填充，
    防止背景图无法完全覆盖时底色透出。
    """
    from pptx.util import Emu
    from lxml import etree
    from pptx.oxml.ns import qn as _qn

    # ─── 清除上游 Skill 设置的 <p:bg> 纯色/渐变填充 ───
    cSld = slide._element.find(_qn('p:cSld'))
    if cSld is not None:
        bg_elem = cSld.find(_qn('p:bg'))
        if bg_elem is not None:
            cSld.remove(bg_elem)

    pres = slide.part.package.presentation_part.presentation
    sw, sh = pres.slide_width, pres.slide_height

    pic = slide.shapes.add_picture(image_path, 0, 0, width=sw, height=sh)

    spTree = slide.shapes._spTree
    pic_el = pic._element
    spTree.remove(pic_el)
    first_shape = None
    for child in spTree.iterchildren():
        tag = etree.QName(child).localname
        if tag in ("sp", "pic", "grpSp", "graphicFrame", "cxnSp"):
            first_shape = child
            break
    if first_shape is not None:
        first_shape.addprevious(pic_el)
    else:
        spTree.append(pic_el)


# ---------- 字体替换（P1-5 修复）----------

def _replace_fonts_in_text_frame(text_frame) -> int:
    """替换单个 text_frame 中所有 run 的字体，返回变更数量。"""
    from lxml import etree as _etree
    changes = 0
    for para in text_frame.paragraphs:
        for run in para.runs:
            size_pt = run.font.size.pt if run.font.size else 18
            is_title = size_pt >= TITLE_SIZE_THRESHOLD_PT

            cur = run.font.name or ""
            if any(k in cur for k in CODE_FONT_KEYWORDS):
                continue

            rPr = run._r.get_or_add_rPr()

            target_latin = FONT_TITLE_LATIN if is_title else FONT_BODY_LATIN
            target_ea = FONT_TITLE_EA if is_title else FONT_BODY_EA

            changed = False

            latin = rPr.find(qn("a:latin"))
            if latin is None:
                latin = _etree.SubElement(rPr, qn("a:latin"))
            if latin.get("typeface") != target_latin:
                latin.set("typeface", target_latin)
                changed = True

            ea = rPr.find(qn("a:ea"))
            if ea is None:
                ea = _etree.SubElement(rPr, qn("a:ea"))
            if ea.get("typeface") != target_ea:
                ea.set("typeface", target_ea)
                changed = True

            cs = rPr.find(qn("a:cs"))
            if cs is None:
                cs = _etree.SubElement(rPr, qn("a:cs"))
            if cs.get("typeface") != target_latin:
                cs.set("typeface", target_latin)
                changed = True

            if changed:
                changes += 1
    return changes


def _replace_fonts_in_shape(shp) -> int:
    """递归处理单个形状的字体替换（支持组合形状和表格）。"""
    changes = 0

    # 组合形状（GroupShape）：递归处理子形状
    if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shp in shp.shapes:
            changes += _replace_fonts_in_shape(child_shp)
        return changes

    # 表格：遍历所有单元格
    if shp.has_table:
        for row in shp.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    changes += _replace_fonts_in_text_frame(cell.text_frame)
        return changes

    # 普通形状：处理 text_frame
    if shp.has_text_frame:
        changes += _replace_fonts_in_text_frame(shp.text_frame)

    return changes


def replace_fonts_in_slide(slide) -> int:
    """遍历所有形状（含组合形状、表格），递归替换字体。

    v6 改进：
    - 递归处理 GroupShape 子形状
    - 处理表格单元格内的文本
    - 只有当目标字体名与当前字体名不同时才计入 changes
    """
    changes = 0
    for shp in slide.shapes:
        changes += _replace_fonts_in_shape(shp)
    return changes


# ---------- Logo（P0-2 去重修复）----------

def _logo_already_exists(slide, left: int, top: int,
                          target_w: int, target_h: int,
                          tolerance: int = 50000) -> bool:
    """检查幻灯片中是否已有同位置/同尺寸的图片形状（Logo 去重）。

    tolerance 单位为 EMU（1pt ≈ 12700 EMU，50000 ≈ 4pt 偏差允许范围）。
    """
    from lxml import etree
    for shp in slide.shapes:
        try:
            sp_el = shp._element
            tag_local = etree.QName(sp_el).localname
            if tag_local != 'pic':
                continue
            # 比较位置和尺寸
            if (abs(shp.left - left) <= tolerance
                    and abs(shp.top - top) <= tolerance
                    and abs(shp.width - target_w) <= tolerance
                    and abs(shp.height - target_h) <= tolerance):
                return True
        except Exception:
            pass
    return False


def check_element_overflow(slide, slide_w: int, slide_h: int,
                            page_type: str) -> list:
    """[P1-B] 检测内容元素是否溢出页面边界或侵入南京页眉安全区。

    返回警告列表（字符串），仅检测不修改。
    """
    from lxml import etree
    warnings = []

    # 南京模板不自动插入横版 Logo，顶部左侧为页眉/标题安全区。
    TEMPLATE_W_INCH = 20.00
    scale = slide_w / Inches(TEMPLATE_W_INCH)
    logo_left = int(Inches(1.20) * scale)
    logo_top = int(Inches(0.35) * scale)
    logo_right = int(Inches(6.00) * scale)
    logo_bottom = int(Inches(1.15) * scale)

    for shp in slide.shapes:
        try:
            sp_el = shp._element
            tag_local = etree.QName(sp_el).localname
            # 跳过背景图片（z-order 底层大图）
            if tag_local == 'pic' and shp.width >= slide_w * 0.9 and shp.height >= slide_h * 0.9:
                continue

            name = getattr(shp, 'name', str(shp.shape_id))

            # 检测溢出页面边界
            right_edge = shp.left + shp.width
            bottom_edge = shp.top + shp.height
            if right_edge > slide_w + Inches(0.1):  # 允许 0.1" 容差
                warnings.append(f"形状 '{name}' 溢出右边界 ({(right_edge - slide_w) / 914400:.2f}\")")
            if bottom_edge > slide_h + Inches(0.1):
                warnings.append(f"形状 '{name}' 溢出下边界 ({(bottom_edge - slide_h) / 914400:.2f}\")")

            # 检测是否侵入南京页眉安全区（只对有文字的形状检测）
            if shp.has_text_frame and shp.text_frame.text.strip():
                shp_right = shp.left + shp.width
                shp_bottom = shp.top + shp.height
                # 矩形相交检测
                if (shp.left < logo_right and shp_right > logo_left and
                        shp.top < logo_bottom and shp_bottom > logo_top):
                    warnings.append(f"形状 '{name}' 侵入南京页眉安全区")
        except Exception:
            pass

    return warnings


def add_logo(slide, page_type: str, logo_main: str, logo_corner: str,
             slide_w: int, slide_h: int) -> bool:
    """南京模板默认不强制插入横版 Logo。

    保留函数作为未来官方南京 Logo 接入钩子；当前永远返回 False。
    """
    return False


# ---------- P0-A：文字色自动反转（对比度修复）----------

def _relative_luminance(r: int, g: int, b: int) -> float:
    """计算相对亮度（WCAG 2.1 定义），范围 [0, 1]。"""
    rl = _rgb_to_linear(r)
    gl = _rgb_to_linear(g)
    bl = _rgb_to_linear(b)
    return 0.2126 * rl + 0.7152 * gl + 0.0722 * bl


def _contrast_ratio(lum1: float, lum2: float) -> float:
    """计算两个亮度值之间的对比度比（WCAG 2.1），范围 [1, 21]。"""
    lighter = max(lum1, lum2)
    darker = min(lum1, lum2)
    return (lighter + 0.05) / (darker + 0.05)


def _get_background_luminance(slide) -> float:
    """估算幻灯片背景亮度。

    策略：
    1. 检查 <p:bg> 中的纯色填充
    2. 若无，按南京浅蓝模板背景估算亮度
    """
    try:
        from lxml import etree
        cSld = slide._element.find(qn('p:cSld'))
        if cSld is not None:
            bg_elem = cSld.find(qn('p:bg'))
            if bg_elem is not None:
                # 尝试找纯色填充
                for srgb in bg_elem.iter(qn('a:srgbClr')):
                    val = srgb.get('val', '')
                    if len(val) == 6:
                        r = int(val[0:2], 16)
                        g = int(val[2:4], 16)
                        b = int(val[4:6], 16)
                        return _relative_luminance(r, g, b)
    except Exception:
        pass
    # 默认：模板使用浅蓝背景（亮度 ≈ 0.92）
    return _relative_luminance(0xDD, 0xEE, 0xFF)


def fix_text_contrast(slide, page_type: str = "content", dry_run: bool = False) -> int:
    """背景替换后检测并修复低对比度文字。

    规则（基于模板为浅色背景的前提）：
    - 所有页面（bg_lum > 0.4）：
        - 封面/章节/结尾页（bg-cover 更浅）：文字亮度 > 0.40 → 改为 #08194B
        - 内容页：文字亮度 > 0.50 且对比度 < 3:1 → 改为 #08194B
    - 背景亮度 ≤ 0.4（深色背景，理论上模板不会出现）：
        - 文字亮度 < 0.2 且对比度 < 3:1 → 改为 #FFFFFF

    仅处理有明确 sRGB 颜色的文字 run，跳过主题色引用。
    返回修复数量。
    """
    bg_lum = _get_background_luminance(slide)
    fixes = 0

    # 封面/章节/结尾页使用更激进的阈值（bg-cover 比 bg-content 更浅）
    is_cover_section_end = page_type in ("cover", "section", "end")
    text_lum_threshold = 0.40 if is_cover_section_end else 0.50

    def _get_shape_fill_luminance(shp) -> float:
        """获取形状自身填充色的亮度，若无填充或非纯色返回 None。"""
        try:
            fill = shp.fill
            rgb = _get_rgb_from_solid_fill(fill)
            if rgb:
                return _relative_luminance(*rgb)
        except Exception:
            pass
        return None

    def _process_runs_in_frame(text_frame, local_bg_lum: float = None):
        """处理 text_frame 中的文字颜色。

        local_bg_lum: 形状自身填充色亮度（若有），优先使用；否则用页面背景亮度。
        """
        nonlocal fixes
        effective_bg_lum = local_bg_lum if local_bg_lum is not None else bg_lum
        # 是否直接在纯背景上（形状无自身填充）
        on_bare_background = (local_bg_lum is None)

        for para in text_frame.paragraphs:
            for run in para.runs:
                # ★ 先检查 run 是否使用 schemeClr 引用白色槽位（不跳过空 run）
                # 这类情况 _get_run_text_color_xml 可能解析不到（主题色被改过）
                # 但实际 PowerPoint 仍然渲染为白色
                if on_bare_background and effective_bg_lum >= 0.35:
                    if _fix_run_scheme_white(run, dry_run):
                        fixes += 1
                        continue

                text = run.text.strip()
                if not text:
                    continue

                orig = _get_run_text_color_xml(run)

                # ★ 关键修复：处理没有显式颜色定义的文字
                # 这类文字继承自 slide layout/master 的 defRPr，
                # 在原深色背景 PPT 中默认显示为白色，迁移到浅色背景后不可见
                if orig is None:
                    if effective_bg_lum >= 0.35:
                        # 浅色/中等背景（页面或形状填充）：无颜色定义的文字强制设为深蓝
                        if not dry_run:
                            _set_run_text_color_xml(run, 0x08, 0x19, 0x4B)
                        fixes += 1
                    else:
                        # 深色背景形状内的无色文字：保持不动（本来就是白色继承，在深底上正确）
                        pass
                    continue

                r, g, b = orig
                text_lum = _relative_luminance(r, g, b)
                cr = _contrast_ratio(effective_bg_lum, text_lum)

                # ★ 浅色背景页面：纯背景（形状无自身填充）上的非黑色文字强制改为黑色
                # 用户明确要求：只有背景的地方，文字一律黑色呈现
                # 适用于所有浅色背景页面（不限 cover/section/end）
                if on_bare_background and effective_bg_lum >= 0.35:
                    # 只要不是已经很深的颜色（亮度<0.10≈接近纯黑），都改为黑色
                    if text_lum > 0.10:
                        if not dry_run:
                            _set_run_text_color_xml(run, 0x08, 0x19, 0x4B)
                        fixes += 1
                    continue

                # 对比度阈值：封面/章节/结尾页更激进
                min_contrast = 2.5 if is_cover_section_end else 3.0

                if cr >= min_contrast:
                    continue  # 对比度足够

                # 浅色/中等背景 + 浅色/白色文字 → 深蓝
                # 阈值 0.35：确保形状填充亮度 ≈ 0.38-0.40 的也能被覆盖
                if effective_bg_lum >= 0.35 and text_lum > text_lum_threshold:
                    if not dry_run:
                        _set_run_text_color_xml(run, 0x08, 0x19, 0x4B)
                    fixes += 1
                # 深色背景 + 深色文字 → 白色
                elif effective_bg_lum < 0.35 and text_lum < 0.2:
                    if not dry_run:
                        _set_run_text_color_xml(run, 0xFF, 0xFF, 0xFF)
                    fixes += 1

    def _process_shape_recursive(shp, depth: int = 0):
        """递归处理形状及其子形状中的文字（支持无限嵌套 GroupShape）。

        如果形状有自身深色填充，文字保留白色（不转黑）。
        """
        if depth > 10:  # 安全上限
            return

        # 检查形状自身填充亮度，用作文字对比度判断的局部背景
        local_bg_lum = _get_shape_fill_luminance(shp)

        if shp.has_text_frame:
            _process_runs_in_frame(shp.text_frame, local_bg_lum)
        if shp.has_table:
            try:
                for row in shp.table.rows:
                    for cell in row.cells:
                        # 表格单元格也可能有自身填充
                        cell_bg_lum = None
                        try:
                            cell_rgb = _get_rgb_from_solid_fill(cell.fill)
                            if cell_rgb:
                                cell_bg_lum = _relative_luminance(*cell_rgb)
                        except Exception:
                            pass
                        if cell.text_frame:
                            _process_runs_in_frame(cell.text_frame, cell_bg_lum)
            except Exception:
                pass
        if hasattr(shp, 'shapes'):
            for sub_shp in shp.shapes:
                _process_shape_recursive(sub_shp, depth + 1)

    for shp in slide.shapes:
        _process_shape_recursive(shp)

    # ★ 额外修复：清除 defRPr / endParaRPr 中的白色/浅色定义
    # 这些不影响已有 run 颜色，但在 PowerPoint 中可能影响空段落、光标色、
    # 以及某些渲染器对无 run 颜色文字的回退判定
    if bg_lum >= 0.35:
        for shp in slide.shapes:
            _fix_defRPr_colors_recursive(shp, dry_run)

    return fixes


def _fix_defRPr_colors_recursive(shp, dry_run: bool = False):
    """递归清除形状中 defRPr / endParaRPr 的白色/浅色颜色定义。

    在浅色背景页面中，段落默认色不应该是白色——否则 PowerPoint 渲染空段落、
    光标颜色或对未设 run 颜色的文字回退时会显示白色。

    处理三层 defRPr：
    1. lstStyle/defRPr —— 形状级列表默认文字色（最容易被遗漏！）
    2. pPr/defRPr    —— 段落级默认文字色
    3. endParaRPr    —— 段尾默认文字色
    """
    # ★ 第1层：形状级 lstStyle/defRPr（PowerPoint 回退最底层，极易遗漏）
    _fix_lstStyle_white(shp, dry_run)

    if shp.has_text_frame:
        _fix_defRPr_in_textframe(shp.text_frame, dry_run)
    if shp.has_table:
        try:
            for row in shp.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        _fix_defRPr_in_textframe(cell.text_frame, dry_run)
        except Exception:
            pass
    if hasattr(shp, 'shapes'):
        for sub_shp in shp.shapes:
            _fix_defRPr_colors_recursive(sub_shp, dry_run)


def _fix_lstStyle_white(shp, dry_run: bool = False):
    """清除形状 txBody/lstStyle 中所有级别（lvl1pPr~lvl9pPr）的 defRPr 白色文字定义。

    OOXML 结构：
      p:txBody / a:lstStyle / a:lvl1pPr / a:defRPr / a:solidFill / a:srgbClr val="FFFFFF"
      p:txBody / a:lstStyle / a:lvl2pPr / a:defRPr / ...
      ...
    以及直接的：
      p:txBody / a:lstStyle / a:defRPr / ...

    这是 PowerPoint 中文字颜色回退的最底层。如果不清除，在浅色背景上文字不可见。

    同时处理 schemeClr 类型的白色引用（bg1/lt1/bg2/lt2）。
    """
    try:
        from lxml import etree as _etree
        # txBody 在 p: 命名空间
        txBody = shp._element.find(qn('p:txBody'))
        if txBody is None:
            return
        lstStyle = txBody.find(qn('a:lstStyle'))
        if lstStyle is None:
            return
        
        # 遍历 lstStyle 的所有子元素（lvl1pPr~lvl9pPr + 直接 defRPr）
        for child in lstStyle:
            local_name = _etree.QName(child).localname
            if local_name == 'defRPr':
                # lstStyle 直接子节点 defRPr
                _replace_light_color_in_rPr_v2(child, dry_run)
            elif local_name.startswith('lvl') and local_name.endswith('pPr'):
                # lvl1pPr ~ lvl9pPr
                defRPr = child.find(qn('a:defRPr'))
                if defRPr is not None:
                    _replace_light_color_in_rPr_v2(defRPr, dry_run)
    except Exception:
        pass


def _fix_defRPr_in_textframe(text_frame, dry_run: bool = False):
    """修复 text_frame 中所有段落的 defRPr 和 endParaRPr 白色颜色。
    
    使用增强版 v2 函数，同时处理 srgbClr 和 schemeClr 类型的白色。
    """
    for para in text_frame.paragraphs:
        pPr = para._p.find(qn('a:pPr'))
        if pPr is not None:
            defRPr = pPr.find(qn('a:defRPr'))
            if defRPr is not None:
                _replace_light_color_in_rPr_v2(defRPr, dry_run)

        endParaRPr = para._p.find(qn('a:endParaRPr'))
        if endParaRPr is not None:
            _replace_light_color_in_rPr_v2(endParaRPr, dry_run)


def _replace_light_color_in_rPr(rPr_elem, dry_run: bool = False):
    """如果 rPr 元素中的 solidFill 是白色/浅色（亮度>0.7），替换为 #08194B。"""
    solidFill = rPr_elem.find('.//' + qn('a:solidFill'))
    if solidFill is None:
        return
    srgbClr = solidFill.find(qn('a:srgbClr'))
    if srgbClr is None:
        return
    val = srgbClr.get('val', '')
    if len(val) != 6:
        return
    r = int(val[0:2], 16)
    g = int(val[2:4], 16)
    b = int(val[4:6], 16)
    lum = _relative_luminance(r, g, b)
    if lum > 0.7:  # 白色或接近白色
        if not dry_run:
            srgbClr.set('val', '08194B')


def _replace_light_color_in_rPr_v2(rPr_elem, dry_run: bool = False):
    """增强版：处理 rPr 元素中的白色/浅色颜色（支持 srgbClr + schemeClr）。

    处理两种情况：
    1. solidFill/srgbClr val="FFFFFF" 等浅色 → 改为 "08194B"
    2. solidFill/schemeClr val="bg1|lt1|bg2|lt2" → 整个替换为 srgbClr val="08194B"
       （因为这些 scheme 槽位在任何主题下都指向浅色/白色）
    """
    from lxml import etree as _etree
    solidFill = rPr_elem.find(qn('a:solidFill'))
    if solidFill is None:
        # 尝试更深层查找
        solidFill = rPr_elem.find('.//' + qn('a:solidFill'))
    if solidFill is None:
        return

    # 情况1：srgbClr 浅色
    srgbClr = solidFill.find(qn('a:srgbClr'))
    if srgbClr is not None:
        val = srgbClr.get('val', '')
        if len(val) == 6:
            r = int(val[0:2], 16)
            g = int(val[2:4], 16)
            b = int(val[4:6], 16)
            lum = _relative_luminance(r, g, b)
            if lum > 0.5:  # 白色或浅色（阈值比v1更激进，覆盖更多浅色）
                if not dry_run:
                    srgbClr.set('val', '08194B')
        return

    # 情况2：schemeClr 指向白色/浅色槽位（bg1/lt1/bg2/lt2）
    schemeClr = solidFill.find(qn('a:schemeClr'))
    if schemeClr is not None:
        slot = schemeClr.get('val', '')
        WHITE_SLOTS = ('bg1', 'lt1', 'bg2', 'lt2')
        if slot in WHITE_SLOTS:
            if not dry_run:
                # 删除 schemeClr，替换为 srgbClr
                solidFill.remove(schemeClr)
                new_srgb = _etree.SubElement(solidFill, qn('a:srgbClr'))
                new_srgb.set('val', '08194B')
        return


# ---------- P2-12：字体安装检测 ----------

def _check_tencent_sans_installed() -> bool:
    """检测 TencentSans 字体是否已安装到系统（跨平台）。"""
    import platform
    system = platform.system()
    try:
        if system == "Windows":
            import winreg
            font_keys = [
                r"SOFTWARE\Microsoft\Windows NT\CurrentVersion\Fonts",
            ]
            for key_path in font_keys:
                try:
                    key = winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, key_path)
                    i = 0
                    while True:
                        try:
                            name, _, _ = winreg.EnumValue(key, i)
                            if "TencentSans" in name or "Tencent Sans" in name:
                                return True
                            i += 1
                        except OSError:
                            break
                    winreg.CloseKey(key)
                except OSError:
                    continue
            # 备用：检查用户字体目录
            user_fonts = Path(os.environ.get("LOCALAPPDATA", "")) / "Microsoft" / "Windows" / "Fonts"
            if user_fonts.exists():
                for f in user_fonts.iterdir():
                    if "TencentSans" in f.name or "Tencent" in f.name:
                        return True
            return False
        elif system == "Darwin":
            # macOS：用 fc-list 或扫描字体目录
            import subprocess
            result = subprocess.run(
                ["fc-list", ":family=TencentSans"],
                capture_output=True, text=True, timeout=3
            )
            return bool(result.stdout.strip())
        else:
            # Linux
            import subprocess
            result = subprocess.run(
                ["fc-list", ":family=TencentSans"],
                capture_output=True, text=True, timeout=3
            )
            return bool(result.stdout.strip())
    except Exception:
        return False  # 检测失败时不阻断流程


# ---------- 主适配函数 ----------

def adapt(input_pptx: str, output_pptx: str,
          city: str = "nanjing",
          mode: str = "full",
          dry_run: bool = False,
          color_compliance: bool = True,
          color_delta: float = 30.0,
          skip_white_black: bool = False) -> dict:
    """模式：
    - full: 替换背景 + 字体 + Logo + 色块合规（默认）
    - light: 仅字体 + Logo + 色块合规，保留原背景
    - logo-only: 仅 Logo

    color_compliance: 是否启用色块合规替换（默认 True）
    color_delta: ΔE 容差阈值，超出阈值不替换（默认 30）
                 注意：禁用色（蓝紫/青/鲜绿）会自动放宽阈值强制替换
    skip_white_black: 不跳过中性色，默认会将黑/灰收敛到南京限定色板
    dry_run=True: 仅分析并输出适配报告，不写入输出文件
    """
    assets = get_assets(city)

    pres = Presentation(input_pptx)
    slide_w = pres.slide_width
    slide_h = pres.slide_height
    total = len(pres.slides)

    # 构建主题色缓存（fix_text_contrast 需要解析 schemeClr 引用）
    _build_theme_color_cache(pres)

    report = {
        "input": input_pptx,
        "output": output_pptx if not dry_run else "(dry-run，不写入文件)",
        "city": city,
        "mode": mode,
        "dry_run": dry_run,
        "color_compliance": color_compliance,
        "color_delta": color_delta,
        "total_slides": total,
        "page_types": [],
        "font_changes": 0,
        "color_changes": [],           # 每条替换记录
        "color_changes_count": 0,
        "theme_color_changes": [],     # 主题色替换记录
        "theme_color_changes_count": 0,
        "overlay_removed": 0,          # [P0-1] 全屏遮罩清除数量
        "logo_skipped": 0,             # [P0-2] Logo/页眉去重跳过数量
        "contrast_fixes": 0,           # [P0-A] 文字对比度修复数量
        "overflow_warnings": [],       # [P1-B] 元素溢出检测
        "warnings": [],
    }

    # [P2-12] 字体安装检测
    if mode in ("full", "light"):
        if not _check_tencent_sans_installed():
            report["warnings"].append(
                "⚠️ TencentSans 字体未安装到系统！字体替换已写入 XML，"
                "但 PowerPoint 打开时可能显示替代字体。"
                f"请手动安装 {ASSETS / 'fonts'} 目录下的 .otf/.ttf 文件后重新打开 PPT。"
            )

    # 0) 主题色覆写（v4 新增：必须在遍历 slide 之前，覆写全局主题定义）
    if color_compliance and mode in ("full", "light"):
        try:
            theme_changes = replace_theme_colors(pres, dry_run=dry_run)
            report["theme_color_changes"] = theme_changes
            report["theme_color_changes_count"] = len(theme_changes)
            report["color_changes"].extend(theme_changes)
            report["color_changes_count"] += len(theme_changes)
        except Exception as e:
            report["warnings"].append(f"主题色覆写失败：{e}")

    for idx, slide in enumerate(pres.slides):
        ptype = classify_page(slide, idx, total)
        report["page_types"].append({"index": idx + 1, "type": ptype})

        # 1) 色块合规替换（在背景替换之前，避免误扫描背景图）
        if color_compliance and mode in ("full", "light"):
            try:
                changes = replace_colors_in_slide(
                    slide,
                    max_delta=color_delta,
                    skip_white_black=skip_white_black,
                    dry_run=dry_run,
                )
                for c in changes:
                    c["slide"] = idx + 1
                report["color_changes"].extend(changes)
                report["color_changes_count"] += len(changes)
            except Exception as e:
                report["warnings"].append(f"第 {idx+1} 页色块合规替换失败：{e}")

        # 2) 背景（仅 full 模式）
        if mode == "full":
            # [P0-1] 插入新背景前，先清除全屏遮罩
            try:
                removed = remove_fullscreen_overlays(
                    slide, slide_w, slide_h, dry_run=dry_run
                )
                report["overlay_removed"] += removed
                if removed > 0:
                    report["warnings"].append(
                        f"第 {idx+1} 页：检测并清除了 {removed} 个全屏遮罩形状（双层背景修复）"
                    )
            except Exception as e:
                report["warnings"].append(f"第 {idx+1} 页全屏遮罩清除失败：{e}")

            try:
                bg = assets.get(f"bg_{ptype}", assets["bg_content"])
                if not dry_run:
                    set_slide_background(slide, bg)
            except Exception as e:
                report["warnings"].append(f"第 {idx+1} 页背景替换失败：{e}")

        # 3) 字体（full / light）
        if mode in ("full", "light"):
            try:
                n = replace_fonts_in_slide(slide)
                report["font_changes"] += n
            except Exception as e:
                report["warnings"].append(f"第 {idx+1} 页字体替换失败：{e}")

        # 3.5) [P0-A] 文字色对比度修复（背景替换后执行）
        if mode in ("full", "light"):
            try:
                contrast_fixed = fix_text_contrast(slide, page_type=ptype, dry_run=dry_run)
                report["contrast_fixes"] += contrast_fixed
            except Exception as e:
                report["warnings"].append(f"第 {idx+1} 页对比度修复失败：{e}")

        # 4) Logo / 页眉：南京模板不再自动插入旧横版 Logo
        try:
            if assets.get("logo_enabled") and assets.get("logo_main") and not dry_run:
                added = add_logo(slide, ptype, assets["logo_main"],
                                 assets.get("logo_corner"), slide_w, slide_h)
                if not added:
                    report["logo_skipped"] += 1
        except Exception as e:
            report["warnings"].append(f"第 {idx+1} 页 Logo/页眉处理失败：{e}")

        # 5) [P1-B] 元素溢出检测
        try:
            overflow_warns = check_element_overflow(slide, slide_w, slide_h, ptype)
            for w in overflow_warns:
                report["overflow_warnings"].append(f"第 {idx+1} 页：{w}")
        except Exception:
            pass

    if not dry_run:
        pres.save(output_pptx)
    else:
        report["warnings"].append("⚠️ dry-run 模式：以上分析仅供预览，文件未写入。")
    return report


# ---------- 报告打印 ----------

def print_report(report: dict):
    print("=" * 70)
    tag = "（DRY-RUN 预览）" if report.get("dry_run") else ""
    print(f"公司沙龙 PPT 模板适配报告 v7 {tag}")
    print("=" * 70)
    print(f"输入：{report['input']}")
    print(f"输出：{report['output']}")
    print(f"城市模板：{report['city']}")
    print(f"适配模式：{report['mode']}")
    print(f"色块合规：{'开启（ΔE < ' + str(report['color_delta']) + '；禁用色强制替换）' if report.get('color_compliance') else '关闭'}")
    print(f"色差算法：CIEDE2000（P1-4 升级）")
    print(f"总页数：{report['total_slides']}")
    print(f"字体替换次数（实际变更）：{report['font_changes']}")
    print(f"主题色覆写次数：{report.get('theme_color_changes_count', 0)}")
    print(f"色块颜色替换次数（含主题）：{report['color_changes_count']}")
    print(f"全屏遮罩清除数量：{report.get('overlay_removed', 0)}")
    print(f"Logo/页眉去重跳过数量：{report.get('logo_skipped', 0)}")
    print(f"文字对比度修复数量：{report.get('contrast_fixes', 0)}")

    if report.get('theme_color_changes'):
        print()
        print("🎨 主题色覆写明细：")
        for c in report['theme_color_changes']:
            print(f"  {c['element']:<20} | {c['from']} → {c['to']} ({c['brand_name']}, ΔE={c['delta_e']})")
    print()
    print("各页适配类型：")
    for p in report["page_types"]:
        print(f"  第 {p['index']:>3} 页 → {p['type']}")

    if report["color_changes_count"] > 0:
        print()
        print("🎨 色块颜色替换明细（前80条）：")
        for i, c in enumerate(report["color_changes"][:80]):
            slide_label = c.get('slide', 0)
            slide_str = f"主题" if slide_label == 0 else f"第{slide_label:>2}页"
            print(f"  {slide_str} | {c['element']:<40} | {c['from']} → {c['to']} ({c['brand_name']}, ΔE={c['delta_e']})")
        if report["color_changes_count"] > 80:
            print(f"  ... 共 {report['color_changes_count']} 条（仅展示前 80 条）")

    if report.get("overflow_warnings"):
        print()
        print("📐 元素溢出/安全区侵入检测：")
        for w in report["overflow_warnings"][:20]:
            print(f"  - {w}")
        if len(report["overflow_warnings"]) > 20:
            print(f"  ... 共 {len(report['overflow_warnings'])} 条（仅展示前 20 条）")

    if report["warnings"]:
        print()
        print("⚠️ 警告/提示：")
        for w in report["warnings"]:
            print(f"  - {w}")
    else:
        print()
        print("✅ 无警告")
    print("=" * 70)


# ---------- CLI ----------

# ---------- 输入格式识别 ----------

# 支持的输入格式及其扩展名
SUPPORTED_PPTX_EXTS = ('.pptx',)
SUPPORTED_HTML_EXTS = ('.html', '.htm', '.xhtml')
SUPPORTED_MD_EXTS = ('.md', '.markdown', '.mdown')

def _detect_input_format(input_path: str) -> str:
    """根据文件扩展名自动检测输入格式。

    返回值: 'pptx' | 'html' | 'markdown' | 'unknown'
    """
    ext = os.path.splitext(input_path)[1].lower()
    if ext in SUPPORTED_PPTX_EXTS:
        return 'pptx'
    elif ext in SUPPORTED_HTML_EXTS:
        return 'html'
    elif ext in SUPPORTED_MD_EXTS:
        return 'markdown'
    return 'unknown'


def _run_html_pipeline(args):
    """HTML/Markdown → PPT 管线：解析 HTML 结构并生成品牌合规 PPT。

    从同目录的 html_to_pptx 模块导入核心类，实现 HTML 文件转换。
    也支持 Markdown 文件（先转 HTML 再走 HTML 管线）。
    """
    # 动态导入 html_to_pptx 模块（与本脚本同目录）
    import importlib.util
    html_module_path = SCRIPT_DIR / "html_to_pptx.py"
    if not html_module_path.exists():
        print(f"[ERROR] html_to_pptx.py 模块不存在: {html_module_path}", file=sys.stderr)
        sys.exit(1)

    spec = importlib.util.spec_from_file_location("html_to_pptx", str(html_module_path))
    html_mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(html_mod)

    input_path = Path(args.input)
    input_format = _detect_input_format(args.input)

    # 读取内容
    encoding = getattr(args, 'encoding', 'utf-8')
    raw_content = input_path.read_text(encoding=encoding)

    # Markdown → HTML 预处理
    if input_format == 'markdown':
        try:
            import markdown as md_lib
            html_content = md_lib.markdown(raw_content, extensions=['tables', 'fenced_code', 'codehilite'])
        except ImportError:
            # 简易 fallback：包裹在 <pre> 中
            print("[WARN] markdown 库未安装，使用简易转换。建议: pip install markdown")
            html_content = f"<html><body><pre>{raw_content}</pre></body></html>"
    else:
        html_content = raw_content

    print(f"[INFO] 检测到输入格式: {input_format.upper()}")
    print(f"[INFO] 解析内容: {args.input}")

    # 检测 HTML 类型并选择解析策略
    if html_mod.is_slide_deck_html(html_content):
        # 幻灯片型 HTML
        print(f"[INFO] 检测到幻灯片型 HTML（slide deck），按页面结构解析")
        slide_contents = html_mod.parse_html_slides(html_content)
        print(f"[INFO] 解析得到 {len(slide_contents)} 个幻灯片页面")
        
        gen = html_mod.PptxGenerator()
        gen.generate_from_slides(slide_contents)
    else:
        # 文章型 HTML
        print(f"[INFO] 检测到文章型 HTML，按内容结构解析")
        blocks = html_mod.parse_html(html_content)
        print(f"[INFO] 解析得到 {len(blocks)} 个内容块")

        # 提取标题
        title = getattr(args, 'title', '') or ''
        if not title:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html_content, 'lxml')
            title_tag = soup.find('title')
            if title_tag:
                title = title_tag.get_text(strip=True)
            elif blocks and blocks[0].type == 'h1':
                title = blocks[0].text
                blocks = blocks[1:]

        subtitle = getattr(args, 'subtitle', '') or ''
        author = getattr(args, 'author', '') or ''

        # 生成 PPT
        print(f"[INFO] 生成 PPT...")
        gen = html_mod.PptxGenerator()
        gen.generate(blocks, title=title, subtitle=subtitle, author=author)

    # 保存
    out_dir = os.path.dirname(os.path.abspath(args.output))
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    gen.save(args.output)

    # 报告
    print(f"\n{'='*60}")
    print(f"[适配报告 - HTML/Markdown → PPT]")
    print(f"{'='*60}")
    print(f"  输入格式: {input_format.upper()}")
    print(f"  输入文件: {args.input}")
    print(f"  输出文件: {args.output}")
    print(f"  总页数:   {gen.slide_count}")
    # 根据解析模式打印不同信息
    if html_mod.is_slide_deck_html(html_content):
        print(f"  模式:     幻灯片型（保留原始 slide 结构）")
        print(f"  页面数:   {len(slide_contents)}")
    else:
        print(f"  标题:     {title}")
        print(f"  内容块:   {len(blocks)}")
    print(f"  品牌合规: 字体=TencentSans W7/W3 | 配色=品牌安全色 | 背景=模板背景 | Logo=南京模板不强制插入")
    print(f"{'='*60}")


def main():
    ap = argparse.ArgumentParser(
        description="腾讯云架构师南京城市沙龙 PPT 模板适配器 v8 — 自动识别输入格式",
        epilog="支持的输入格式: .pptx(迁移适配) | .html/.htm(HTML转PPT) | .md(Markdown转PPT)"
    )
    ap.add_argument("--input", required=True,
                    help="输入文件路径（支持 .pptx / .html / .htm / .md，自动识别格式）")
    ap.add_argument("--output", required=True, help="输出 .pptx 路径")

    # PPTX 迁移模式参数
    ap.add_argument("--city", default="nanjing", help="城市模板（默认 nanjing）[仅 PPTX 模式]")
    ap.add_argument("--mode", default="full", choices=["full", "light", "logo-only"],
                    help="适配模式：full=背景+字体+Logo（默认）；light=字体+Logo；logo-only=仅Logo [仅 PPTX 模式]")
    ap.add_argument("--dry-run", action="store_true",
                    help="仅分析并打印适配报告，不写入输出文件 [仅 PPTX 模式]")
    ap.add_argument("--color-compliance", action="store_true", default=True,
                    help="启用色块合规替换（默认开启）[仅 PPTX 模式]")
    ap.add_argument("--no-color-compliance", action="store_false", dest="color_compliance",
                    help="关闭色块合规替换 [仅 PPTX 模式]")
    ap.add_argument("--color-delta", type=float, default=30.0,
                    help="颜色替换 ΔE 容差阈值（默认 30）[仅 PPTX 模式]")
    ap.add_argument("--no-skip-white-black", action="store_false", dest="skip_white_black",
                    help="兼容旧参数：南京模式默认已不跳过中性色 [仅 PPTX 模式]")
    ap.set_defaults(skip_white_black=False)

    # HTML/Markdown 模式参数
    ap.add_argument("--title", "-t", default="", help="PPT 主标题（默认从内容提取）[仅 HTML/MD 模式]")
    ap.add_argument("--subtitle", "-s", default="", help="PPT 副标题 [仅 HTML/MD 模式]")
    ap.add_argument("--author", "-a", default="", help="作者/演讲者 [仅 HTML/MD 模式]")
    ap.add_argument("--encoding", default="utf-8", help="输入文件编码（默认 utf-8）[仅 HTML/MD 模式]")

    args = ap.parse_args()

    # ---------- 检查输入文件 ----------
    if not os.path.exists(args.input):
        print(f"[ERROR] 输入文件不存在：{args.input}", file=sys.stderr)
        sys.exit(1)

    # ---------- 自动路由 ----------
    input_format = _detect_input_format(args.input)

    if input_format == 'pptx':
        # PPTX → PPTX 迁移适配
        print(f"[INFO] 检测到输入格式: PPTX → 走迁移适配管线")
        if not args.dry_run:
            out_dir = os.path.dirname(os.path.abspath(args.output))
            if out_dir:
                os.makedirs(out_dir, exist_ok=True)

        report = adapt(
            args.input,
            args.output,
            city=args.city,
            mode=args.mode,
            dry_run=args.dry_run,
            color_compliance=args.color_compliance,
            color_delta=args.color_delta,
            skip_white_black=args.skip_white_black,
        )
        print_report(report)

    elif input_format in ('html', 'markdown'):
        # HTML/Markdown → PPT 直接生成
        print(f"[INFO] 检测到输入格式: {input_format.upper()} → 走 HTML 生成管线")
        _run_html_pipeline(args)

    else:
        ext = os.path.splitext(args.input)[1]
        print(f"[ERROR] 不支持的输入格式: '{ext}'", file=sys.stderr)
        print(f"        支持的格式: .pptx | .html/.htm | .md/.markdown", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
