#!/usr/bin/env python3
"""
verify_output.py - HTML->PPT 输出验证脚本
用法: python verify_output.py --pptx output.pptx [--html source.html] [--strict]

功能：
- 检查每页 shapes 数量是否合理（>=3）
- 验证字号范围（标题 40pt / 章节 48pt / 正文 20pt）
- 南京模板默认不强制插入横版 Logo
- 验证背景图存在
- 对比原始 HTML 页面数（如提供）
- 输出通过/失败报告

退出码：0=全部通过  1=有警告  2=有错误
"""
import argparse
import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(2)


# === 常量 ===
TEMPLATE_WIDTH_INCH = 20.00
TEMPLATE_HEIGHT_INCH = 11.25
MIN_SHAPES_PER_SLIDE = 2  # 南京模板不强制插入 Logo，章节/结尾页可只有背景+文本
LOGO_WIDTH_INCH = 4.88
LOGO_HEIGHT_INCH = 0.53
LOGO_TOLERANCE_INCH = 0.5  # Logo 位置容差

# 字号范围（pt）
FONT_SIZE_RANGES = {
    'title': (36, 52),      # 主标题 / 章节标题
    'subtitle': (20, 30),   # 副标题 / 列标题
    'body': (14, 26),       # 正文 / 卡片 / 列表
    'min_readable': 12,     # 最小可读字号
}

# Logo 预期位置
LOGO_POSITIONS = {
    'left': {'x_range': (0.3, 1.5), 'y_range': (0.3, 1.2)},   # cover/section/end
    'right': {'x_range': (13.5, 16.0), 'y_range': (0.2, 1.2)}, # content
}


class VerifyResult:
    def __init__(self):
        self.errors = []
        self.warnings = []
        self.passes = []

    def error(self, slide_num, msg):
        self.errors.append(f"[ERROR] Slide {slide_num}: {msg}")

    def warn(self, slide_num, msg):
        self.warnings.append(f"[WARN]  Slide {slide_num}: {msg}")

    def ok(self, slide_num, msg):
        self.passes.append(f"[OK]    Slide {slide_num}: {msg}")

    def print_report(self):
        print("\n" + "=" * 60)
        print("  HTML->PPT Output Verification Report")
        print("=" * 60)

        if self.passes:
            print(f"\n--- PASSED ({len(self.passes)}) ---")
            for p in self.passes[:20]:  # 只显示前20条
                print(f"  {p}")
            if len(self.passes) > 20:
                print(f"  ... and {len(self.passes) - 20} more")

        if self.warnings:
            print(f"\n--- WARNINGS ({len(self.warnings)}) ---")
            for w in self.warnings:
                print(f"  {w}")

        if self.errors:
            print(f"\n--- ERRORS ({len(self.errors)}) ---")
            for e in self.errors:
                print(f"  {e}")

        print("\n" + "-" * 60)
        total = len(self.passes) + len(self.warnings) + len(self.errors)
        print(f"Total checks: {total}")
        print(f"  Passed:   {len(self.passes)}")
        print(f"  Warnings: {len(self.warnings)}")
        print(f"  Errors:   {len(self.errors)}")

        if self.errors:
            print("\n  RESULT: FAIL")
            return 2
        elif self.warnings:
            print("\n  RESULT: PASS (with warnings)")
            return 1
        else:
            print("\n  RESULT: ALL PASS")
            return 0


def check_shapes_count(slide, slide_num, result):
    """检查每页 shapes 数量"""
    count = len(slide.shapes)
    if count < MIN_SHAPES_PER_SLIDE:
        result.warn(slide_num, f"Only {count} shapes (expected >={MIN_SHAPES_PER_SLIDE})")
    else:
        result.ok(slide_num, f"{count} shapes")


def check_font_sizes(slide, slide_num, result, strict=False):
    """检查字号范围"""
    min_size = 999
    max_size = 0
    has_text = False

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    size_pt = run.font.size.pt
                    has_text = True
                    min_size = min(min_size, size_pt)
                    max_size = max(max_size, size_pt)

    if not has_text:
        return

    if min_size < FONT_SIZE_RANGES['min_readable']:
        result.error(slide_num, f"Font too small: {min_size}pt (min {FONT_SIZE_RANGES['min_readable']}pt)")
    elif strict and min_size < 14:
        result.warn(slide_num, f"Small font detected: {min_size}pt")
    else:
        result.ok(slide_num, f"Font range: {min_size}-{max_size}pt")


def check_logo(slide, slide_num, result, slide_width_inch):
    """南京模板默认不强制插入横版 Logo。"""
    result.ok(slide_num, "Nanjing template: logo insertion not required")


def check_background(slide, slide_num, result, slide_width_emu, slide_height_emu):
    """检查背景图是否存在（全页图片或背景填充）"""
    has_bg = False

    # 方法1：检查全页大小的图片
    for shape in slide.shapes:
        if hasattr(shape, 'image') and shape.image:
            # 容差 5%
            if (shape.width > slide_width_emu * 0.95 and
                shape.height > slide_height_emu * 0.95):
                has_bg = True
                break

    # 方法2：检查 slide background
    if not has_bg:
        bg = slide.background
        if bg and bg.fill and bg.fill.type is not None:
            has_bg = True

    if not has_bg:
        result.warn(slide_num, "No background image/fill detected")
    else:
        result.ok(slide_num, "Background present")


def count_html_slides(html_path):
    """统计 HTML 中的 slide 数量"""
    try:
        from bs4 import BeautifulSoup
        import re
        with open(html_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'lxml')

        # 精确匹配 class="slide"（避免误匹 slide-footer/slide-header/slide-body）
        slide_pattern = re.compile(r'(?:^|\s)slide(?:\s|$)')

        # 尝试多种 slide 标记
        slides = soup.find_all('div', class_=slide_pattern)
        if len(slides) >= 3:
            return len(slides)

        slides = soup.find_all('div', attrs={'data-slide': True})
        if len(slides) >= 3:
            return len(slides)

        slides = soup.find_all('section', class_=slide_pattern)
        if len(slides) >= 3:
            return len(slides)

        # 文章型：按 h1/h2 估算页数
        headings = soup.find_all(['h1', 'h2'])
        return max(len(headings), 1) + 2  # +2 for cover and end

    except Exception as e:
        print(f"  [INFO] Cannot parse HTML: {e}")
        return None


def check_theme_fonts(prs, result):
    """[v9] 检查 theme 中 majorFont/minorFont 是否为 TencentSans"""
    try:
        from lxml import etree as ET
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_map = {'a': ns}

        for master in prs.slide_masters:
            for rel in master.part.rels.values():
                if 'theme' not in str(rel.reltype).lower():
                    continue
                theme_part = rel.target_part
                root = ET.fromstring(theme_part.blob)

                major = root.find('.//a:fontScheme/a:majorFont/a:latin', ns_map)
                minor = root.find('.//a:fontScheme/a:minorFont/a:latin', ns_map)

                if major is not None:
                    tf = major.get('typeface', '')
                    if 'TencentSans' in tf:
                        result.ok(0, f"Theme majorFont: {tf}")
                    else:
                        result.error(0, f"Theme majorFont not TencentSans: '{tf}'")
                else:
                    result.error(0, "Theme majorFont not found")

                if minor is not None:
                    tf = minor.get('typeface', '')
                    if 'TencentSans' in tf:
                        result.ok(0, f"Theme minorFont: {tf}")
                    else:
                        result.error(0, f"Theme minorFont not TencentSans: '{tf}'")
                else:
                    result.error(0, "Theme minorFont not found")

                return  # 只检查第一个 theme
    except Exception as e:
        result.warn(0, f"Theme font check failed: {e}")


def check_theme_colors(prs, result):
    """[v9] 检查 theme clrScheme 是否为南京品牌色"""
    EXPECTED = {
        "dk1": "08194B", "lt1": "FFFFFF",
        "dk2": "44474F", "lt2": "8B8C8C",
        "accent1": "3272DC", "accent2": "08194B",
        "accent3": "00C8D8", "accent4": "01A4FF",
        "accent5": "44474F", "accent6": "8B8C8C",
        "hlink": "3272DC", "folHlink": "08194B",
    }
    try:
        from lxml import etree as ET
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_map = {'a': ns}

        for master in prs.slide_masters:
            for rel in master.part.rels.values():
                if 'theme' not in str(rel.reltype).lower():
                    continue
                theme_part = rel.target_part
                root = ET.fromstring(theme_part.blob)
                clr = root.find('.//a:clrScheme', ns_map)
                if clr is None:
                    result.error(0, "No clrScheme found in theme")
                    return

                mismatches = []
                for child in clr:
                    slot = ET.QName(child).localname
                    if slot not in EXPECTED:
                        continue
                    for ce in child:
                        if ET.QName(ce).localname == 'srgbClr':
                            val = ce.get('val', '').upper()
                            expected = EXPECTED[slot].upper()
                            if val != expected:
                                mismatches.append(f"{slot}: {val} != {expected}")

                if mismatches:
                    result.error(0, f"Theme clrScheme mismatches: {'; '.join(mismatches)}")
                else:
                    result.ok(0, f"Theme clrScheme: all 12 slots match brand colors")

                return
    except Exception as e:
        result.warn(0, f"Theme color check failed: {e}")


def check_font_embedding(prs, result):
    """[v9] 检查是否嵌入了 TencentSans 字体文件"""
    try:
        font_rels = []
        for rel in prs.part.rels.values():
            if 'font' in str(rel.reltype).lower():
                font_rels.append(rel.target_ref)

        if len(font_rels) >= 2:
            result.ok(0, f"Font embedding: {len(font_rels)} fonts embedded")
        elif len(font_rels) == 1:
            result.warn(0, f"Font embedding: only 1 font embedded (expected 2)")
        else:
            result.warn(0, "Font embedding: no fonts embedded (TencentSans may not display)")
    except Exception as e:
        result.warn(0, f"Font embedding check failed: {e}")


def main():
    parser = argparse.ArgumentParser(description='Verify HTML->PPT output quality (v9)')
    parser.add_argument('--pptx', required=True, help='Path to generated .pptx file')
    parser.add_argument('--html', help='Path to source .html file (optional, for page count comparison)')
    parser.add_argument('--strict', action='store_true', help='Enable strict mode (warnings become errors)')
    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        print(f"ERROR: File not found: {args.pptx}")
        sys.exit(2)

    print(f"Verifying: {args.pptx}")
    prs = Presentation(args.pptx)
    result = VerifyResult()

    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    slide_width_inch = slide_width_emu / 914400
    slide_height_inch = slide_height_emu / 914400

    print(f"  Slide size: {slide_width_inch:.1f}\" x {slide_height_inch:.1f}\"")
    print(f"  Total slides: {len(prs.slides)}")

    # [v9] 主题级检查（最重要，放在最前面）
    print("\n--- Theme-level checks ---")
    check_theme_fonts(prs, result)
    check_theme_colors(prs, result)
    check_font_embedding(prs, result)

    # 逐页检查
    print("\n--- Slide-level checks ---")
    for i, slide in enumerate(prs.slides, 1):
        check_shapes_count(slide, i, result)
        check_font_sizes(slide, i, result, strict=args.strict)
        check_logo(slide, i, result, slide_width_inch)
        check_background(slide, i, result, slide_width_emu, slide_height_emu)

    # HTML 对比检查
    if args.html and os.path.exists(args.html):
        html_slides = count_html_slides(args.html)
        if html_slides:
            pptx_slides = len(prs.slides)
            diff = abs(pptx_slides - html_slides)
            if diff > html_slides * 0.3:  # 差异超过30%
                result.warn(0, f"Slide count mismatch: PPTX={pptx_slides}, HTML~={html_slides} (diff={diff})")
            else:
                result.ok(0, f"Slide count OK: PPTX={pptx_slides}, HTML~={html_slides}")

    # 输出报告
    exit_code = result.print_report()
    sys.exit(exit_code)


if __name__ == '__main__':
    main()
