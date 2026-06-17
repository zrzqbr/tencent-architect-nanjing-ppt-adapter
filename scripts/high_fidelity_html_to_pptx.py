#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""High-fidelity HTML slide deck to Nanjing template PPTX renderer.

This module is intentionally template-aware but not input-fixed:
- It first analyzes the user's HTML slide deck structure.
- It classifies each slide as cover / section / content / end with generic rules.
- It renders the original HTML components with Chrome to preserve complex cards,
  grids, timelines, product cards, and custom CSS that python-pptx cannot rebuild
  faithfully.
- It injects TencentSans fonts and Nanjing palette overrides.
- It adapts font sizes per slide based on content density.
- It overlays each transparent rendered slide on the matching Nanjing template
  background in PPTX.

Trade-off: this mode optimizes visual fidelity. Text/components are rendered as
images and are not fully editable in PowerPoint. Use html_to_pptx.py editable mode
when editability is more important than fidelity.
"""

from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches

SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent
ASSETS = SKILL_DIR / "assets"
BACKGROUNDS = ASSETS / "backgrounds"
FONTS = ASSETS / "fonts"

SLIDE_CLASS_PATTERN = re.compile(r"(?:^|\s)slide(?:\s|$)")

SLIDE_WIDTH_PX = 1920
SLIDE_HEIGHT_PX = 1080
SLIDE_WIDTH_IN = 20.0
SLIDE_HEIGHT_IN = 11.25

NANJING_COLORS = {
    "primary": "#3272DC",
    "navy": "#08194B",
    "white": "#FFFFFF",
    "cyan": "#00C8D8",
    "bright_blue": "#01A4FF",
    "dark_gray": "#44474F",
    "mid_gray": "#8B8C8C",
}


@dataclass
class SlideAnalysis:
    index: int
    slide_type: str
    text_length: int
    component_count: int
    title: str
    density: str
    font_scale: float


def _has_class_token(element, token: str) -> bool:
    classes = element.get("class", []) or []
    return token in classes


def _find_slides(html_content: str):
    soup = BeautifulSoup(html_content, "lxml")
    slides = soup.find_all("div", class_=SLIDE_CLASS_PATTERN)
    if not slides:
        slides = soup.find_all("div", attrs={"data-slide": True})
    if not slides:
        slides = soup.find_all("section", class_=SLIDE_CLASS_PATTERN)
    return soup, slides


def is_slide_deck_html(html_content: str) -> bool:
    """Return True only when exact slide class tokens are found."""
    _, slides = _find_slides(html_content)
    return len(slides) >= 3


def _text(element, selector: str) -> str:
    found = element.select_one(selector)
    return found.get_text(" ", strip=True) if found else ""


def _component_count(slide) -> int:
    selectors = [
        ".card", ".stat-block", ".product-card", ".scene-card",
        ".question-card", ".sales-step", ".timeline-item", ".tech-list-item",
        ".info-card", ".banner", ".flow-step", ".summary-card",
        ".policy-card", ".highlight-box", "table", "img", "svg", "canvas",
    ]
    return sum(len(slide.select(selector)) for selector in selectors)


def _classify_slide(slide, index: int, total: int) -> str:
    classes = " ".join(slide.get("class", []) or [])
    text = slide.get_text(" ", strip=True)
    title = (
        _text(slide, ".cover-title") or _text(slide, ".section-title") or
        _text(slide, ".section-heading") or _text(slide, "h1") or _text(slide, "h2")
    )
    lower = text.lower()

    if (
        index == 0 or
        "title-page" in classes or "slide-title" in classes or
        slide.select_one(".cover-layout") or slide.select_one(".cover-title") or
        slide.select_one(".cover-eyebrow")
    ):
        return "cover"

    if (
        "slide-end" in classes or
        (index == total - 1 and any(k in text for k in ["感谢", "谢谢", "结语"])) or
        (index == total - 1 and any(k in lower for k in ["thanks", "thank you", "q&a", "qa"]))
    ):
        return "end"

    # Section/agenda pages are transition pages, not every page with a section number.
    # Many real slide decks use .section-num on normal content pages, so combine
    # title semantics with component density instead of hardcoding class presence.
    if "slide-center" in classes or any(k in title for k in ["目录", "议程", "分享结构", "Agenda"]):
        return "section"

    return "content"


def _density_and_scale(slide, slide_type: str) -> tuple[str, float]:
    text_len = len(slide.get_text(" ", strip=True))
    comp = _component_count(slide)

    # Adaptive, not fixed: every slide gets a scale based on information density.
    # Sparse slides can become more expressive; dense slides still receive a slight
    # increase but avoid overflow.
    if slide_type == "cover":
        return "cover", 1.10
    if slide_type == "end":
        return "end", 1.08
    if text_len > 1100 or comp >= 14:
        return "dense", 1.02
    if text_len > 700 or comp >= 8:
        return "rich", 1.05
    if text_len > 350 or comp >= 4:
        return "medium", 1.08
    return "sparse", 1.14


def analyze_html_slides(html_content: str) -> List[SlideAnalysis]:
    _, slides = _find_slides(html_content)
    total = len(slides)
    result: List[SlideAnalysis] = []
    for idx, slide in enumerate(slides):
        slide_type = _classify_slide(slide, idx, total)
        density, scale = _density_and_scale(slide, slide_type)
        title = (
            _text(slide, ".cover-title") or _text(slide, ".section-title") or
            _text(slide, ".section-heading") or _text(slide, "h1") or _text(slide, "h2")
        )
        result.append(
            SlideAnalysis(
                index=idx + 1,
                slide_type=slide_type,
                text_length=len(slide.get_text(" ", strip=True)),
                component_count=_component_count(slide),
                title=title,
                density=density,
                font_scale=scale,
            )
        )
    return result


def _find_chrome(chrome_path: Optional[str] = None) -> str:
    candidates = []
    if chrome_path:
        candidates.append(chrome_path)
    env_path = os.environ.get("CHROME_PATH")
    if env_path:
        candidates.append(env_path)
    candidates.extend([
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
        shutil.which("google-chrome") or "",
        shutil.which("chromium") or "",
        shutil.which("chromium-browser") or "",
    ])
    for item in candidates:
        if item and Path(item).exists():
            return item
    raise RuntimeError("未找到 Chrome/Edge/Chromium。高保真模式需要本机浏览器渲染 HTML。")


def _find_node(node_path: Optional[str] = None) -> str:
    candidates = []
    if node_path:
        candidates.append(node_path)
    env_path = os.environ.get("NODE_BIN")
    if env_path:
        candidates.append(env_path)
    candidates.extend([
        "/Users/raelzhang/.workbuddy/binaries/node/versions/22.22.2/bin/node",
        shutil.which("node") or "",
    ])
    for item in candidates:
        if item and Path(item).exists():
            return item
    raise RuntimeError("未找到 Node.js。高保真模式需要 Node.js + puppeteer-core。")


def _node_env() -> dict:
    env = os.environ.copy()
    node_paths = [
        str(SKILL_DIR / "node_modules"),
        str(SKILL_DIR / ".." / "node_modules"),
        "/Users/raelzhang/.workbuddy/binaries/node/workspace/node_modules",
    ]
    existing = [p for p in node_paths if Path(p).exists()]
    if existing:
        current = env.get("NODE_PATH", "")
        env["NODE_PATH"] = os.pathsep.join(existing + ([current] if current else []))
    return env


def _file_url(path: Path) -> str:
    return path.resolve().as_uri()


def _build_render_js(input_html: Path, out_dir: Path, analysis: List[SlideAnalysis], chrome_path: str) -> str:
    font_w3 = FONTS / "TencentSans-W3.ttf"
    font_w7 = FONTS / "TencentSans-W7.ttf"
    if not font_w3.exists() or not font_w7.exists():
        raise RuntimeError("缺少 TencentSans 字体文件：assets/fonts/TencentSans-W3.ttf 和 TencentSans-W7.ttf")

    payload = {
        "inputUrl": _file_url(input_html),
        "outDir": str(out_dir),
        "chromePath": chrome_path,
        "fontW3": _file_url(font_w3),
        "fontW7": _file_url(font_w7),
        "analysis": [a.__dict__ for a in analysis],
        "colors": NANJING_COLORS,
        "width": SLIDE_WIDTH_PX,
        "height": SLIDE_HEIGHT_PX,
    }

    return r'''
const puppeteer = require('puppeteer-core');
const fs = require('fs');
const path = require('path');
const payload = __PAYLOAD__;
fs.mkdirSync(payload.outDir, { recursive: true });

function css(colors, fontW3, fontW7) {
  return `
@font-face { font-family: 'TencentSans W3'; src: url('${fontW3}') format('truetype'); font-weight: 400; }
@font-face { font-family: 'TencentSans W7'; src: url('${fontW7}') format('truetype'); font-weight: 700; }
:root {
  --bg: transparent !important;
  --bg2: transparent !important;
  --primary: ${colors.primary} !important;
  --primary2: ${colors.bright_blue} !important;
  --accent: ${colors.cyan} !important;
  --accent2: ${colors.bright_blue} !important;
  --text: ${colors.navy} !important;
  --text-dim: ${colors.dark_gray} !important;
  --gold: ${colors.primary} !important;
  --border: rgba(50,114,220,0.28) !important;
  --glass: rgba(255,255,255,0.88) !important;
}
html, body { background: transparent !important; font-family: 'TencentSans W3', 'PingFang SC', 'Microsoft YaHei', sans-serif !important; }
.slide { background: transparent !important; background-image: none !important; }
.bg-grid, .bg-radial, .bg-radial-accent, .bg-radial-green { background: transparent !important; background-image: none !important; }
.orb { display: none !important; }
* { font-family: 'TencentSans W3', 'PingFang SC', 'Microsoft YaHei', sans-serif !important; }
.cover-title, .section-title, .card-title, .product-name, .scene-title, .question-text, .timeline-title, .tech-list-title, .banner-title, .info-title, .sales-step-title, .method-title, .col-title {
  font-family: 'TencentSans W7', 'PingFang SC', 'Microsoft YaHei', sans-serif !important;
  color: ${colors.navy} !important;
  -webkit-text-fill-color: ${colors.navy} !important;
  background: none !important;
}
.cover-title { color:${colors.navy} !important; -webkit-text-fill-color:${colors.navy} !important; }
.cover-sub, .cover-eyebrow, .meta-value, .section-num, .section-tag, .header-page, .footer-right, .stat-num, .card-metric, .spec-val, .tag, .badge, .question-label, .timeline-year, .tech-list-num {
  color: ${colors.primary} !important;
  -webkit-text-fill-color: ${colors.primary} !important;
  background: none !important;
}
.cover-desc, .card-text, .product-desc, .scene-desc, .question-answer, .tech-list-desc, .timeline-desc, .banner-desc, .section-subtitle, .header-title, .footer-left, .meta-label, .stat-label, .stat-desc, .method-desc, .method-analogy, .sales-step-desc, .col-text {
  color: ${colors.dark_gray} !important;
  -webkit-text-fill-color: ${colors.dark_gray} !important;
}
.card, .product-card, .scene-card, .question-card, .info-card, .banner, .tech-list-item, .flow-item, .flow-step, .policy-card, .summary-card, .highlight-box, .stat-block {
  background: rgba(255,255,255,0.88) !important;
  border: 1px solid rgba(50,114,220,0.28) !important;
  box-shadow: 0 12px 32px rgba(8,25,75,0.08) !important;
  backdrop-filter: blur(8px) !important;
}
.card::before { background: linear-gradient(90deg, transparent, ${colors.primary}, transparent) !important; }
.product-card::after { background: linear-gradient(90deg, ${colors.primary}, ${colors.cyan}) !important; }
.header-logo, .tech-center, .sales-step-num, .tech-list-num {
  background: linear-gradient(135deg, ${colors.primary}, ${colors.bright_blue}) !important;
  color: ${colors.white} !important;
  -webkit-text-fill-color: ${colors.white} !important;
}
.tech-ring { border-color: rgba(50,114,220,0.28) !important; }
.tech-node { background: ${colors.cyan} !important; box-shadow: 0 0 15px ${colors.cyan} !important; }
.tag, .badge-primary, .badge-green, .badge-purple, .badge-gold, .badge-red {
  background: rgba(50,114,220,0.10) !important;
  border: 1px solid rgba(50,114,220,0.22) !important;
  color: ${colors.primary} !important;
  -webkit-text-fill-color: ${colors.primary} !important;
}
.slide-header { padding-top: 26px !important; }
.slide-body { padding-bottom: 34px !important; }
.footer-dot, .dot { background:${colors.cyan} !important; }
svg, canvas { max-width: 100% !important; }
`;
}

function applyAdaptiveFonts(slides, analysis) {
  slides.forEach((slide, idx) => {
    const meta = analysis[idx] || { font_scale: 1.06, density: 'medium', slide_type: 'content' };
    slide.setAttribute('data-density', meta.density);
    slide.setAttribute('data-slide-type', meta.slide_type);
    const scale = meta.font_scale || 1.06;
    const elements = slide.querySelectorAll('*');
    elements.forEach(el => {
      const style = window.getComputedStyle(el);
      if (!style || style.display === 'none' || style.visibility === 'hidden') return;
      const fs = parseFloat(style.fontSize || '0');
      if (!fs || fs < 8 || fs > 96) return;
      let factor = scale;
      if (el.closest('.cover-title')) factor = Math.min(scale + 0.04, 1.18);
      if (el.closest('.slide-footer') || el.closest('.footer-right')) factor = Math.min(scale, 1.06);
      const next = Math.max(11, Math.min(90, fs * factor));
      el.style.fontSize = `${next}px`;
    });
  });
}

(async () => {
  const browser = await puppeteer.launch({
    executablePath: payload.chromePath,
    headless: 'new',
    args: ['--allow-file-access-from-files', '--no-sandbox', '--disable-web-security']
  });
  const page = await browser.newPage();
  await page.setViewport({ width: payload.width, height: payload.height, deviceScaleFactor: 2 });
  await page.goto(payload.inputUrl, { waitUntil: 'networkidle0', timeout: 60000 });
  await page.addStyleTag({ content: css(payload.colors, payload.fontW3, payload.fontW7) });
  await page.evaluate((analysis) => {
    function applyAdaptiveFontsInPage(slides, analysisData) {
      slides.forEach((slide, idx) => {
        const meta = analysisData[idx] || { font_scale: 1.06, density: 'medium', slide_type: 'content' };
        slide.setAttribute('data-density', meta.density);
        slide.setAttribute('data-slide-type', meta.slide_type);
        const scale = meta.font_scale || 1.06;
        const elements = slide.querySelectorAll('*');
        elements.forEach(el => {
          const style = window.getComputedStyle(el);
          if (!style || style.display === 'none' || style.visibility === 'hidden') return;
          const fs = parseFloat(style.fontSize || '0');
          if (!fs || fs < 8 || fs > 96) return;
          let factor = scale;
          if (el.closest('.cover-title')) factor = Math.min(scale + 0.04, 1.18);
          if (el.closest('.slide-footer') || el.closest('.footer-right')) factor = Math.min(scale, 1.06);
          const next = Math.max(11, Math.min(90, fs * factor));
          el.style.fontSize = `${next}px`;
        });
      });
    }
    const slides = Array.from(document.querySelectorAll('.slide'));
    applyAdaptiveFontsInPage(slides, analysis);
  }, payload.analysis);
  await page.evaluateHandle('document.fonts.ready');
  await new Promise(resolve => setTimeout(resolve, 500));

  const slides = await page.$$('.slide');
  console.log(`slides=${slides.length}`);
  for (let i = 0; i < slides.length; i++) {
    const slide = slides[i];
    await slide.evaluate(el => el.scrollIntoView({ block: 'start', inline: 'start' }));
    await new Promise(resolve => setTimeout(resolve, 120));
    const out = path.join(payload.outDir, `slide_${String(i + 1).padStart(2, '0')}.png`);
    await slide.screenshot({ path: out, omitBackground: true });
    console.log(out);
  }
  await browser.close();
})().catch(err => {
  console.error(err && err.stack ? err.stack : String(err));
  process.exit(1);
});
'''.replace("__PAYLOAD__", json.dumps(payload, ensure_ascii=False))


def _background_for_slide(slide_type: str) -> Path:
    mapping = {
        "cover": BACKGROUNDS / "bg-cover.jpeg",
        "section": BACKGROUNDS / "bg-section.jpeg",
        "content": BACKGROUNDS / "bg-content.jpeg",
        "end": BACKGROUNDS / "bg-end.jpeg",
    }
    return mapping.get(slide_type, BACKGROUNDS / "bg-content.jpeg")


def _build_pptx(rendered_dir: Path, output_pptx: Path, analysis: List[SlideAnalysis]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank = prs.slide_layouts[6]

    images = sorted(rendered_dir.glob("slide_*.png"))
    if len(images) != len(analysis):
        raise RuntimeError(f"渲染页数与分析页数不一致：rendered={len(images)}, analysis={len(analysis)}")

    for img, meta in zip(images, analysis):
        slide = prs.slides.add_slide(blank)
        bg = _background_for_slide(meta.slide_type)
        if not bg.exists():
            raise RuntimeError(f"背景文件缺失：{bg}")
        slide.shapes.add_picture(str(bg), 0, 0, width=prs.slide_width, height=prs.slide_height)
        slide.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=prs.slide_height)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))


def render_html_to_pptx(
    input_html: str,
    output_pptx: str,
    chrome_path: Optional[str] = None,
    node_path: Optional[str] = None,
    keep_images: bool = False,
) -> dict:
    input_path = Path(input_html)
    output_path = Path(output_pptx)
    html_content = input_path.read_text(encoding="utf-8")
    analysis = analyze_html_slides(html_content)
    if not analysis:
        raise RuntimeError("未识别到幻灯片型 HTML。")

    chrome = _find_chrome(chrome_path)
    node = _find_node(node_path)

    tmp_root_obj = tempfile.TemporaryDirectory(prefix="nanjing_html_render_")
    tmp_root = Path(tmp_root_obj.name)
    rendered_dir = tmp_root / "slides"
    js_path = tmp_root / "render.js"
    rendered_dir.mkdir(parents=True, exist_ok=True)
    js_path.write_text(_build_render_js(input_path, rendered_dir, analysis, chrome), encoding="utf-8")

    try:
        subprocess.run(
            [node, str(js_path)],
            check=True,
            cwd=str(SKILL_DIR),
            env=_node_env(),
            text=True,
        )
        _build_pptx(rendered_dir, output_path, analysis)
        if keep_images:
            keep_dir = output_path.with_suffix("")
            keep_dir = keep_dir.parent / f"{keep_dir.name}-rendered-slides"
            if keep_dir.exists():
                shutil.rmtree(keep_dir)
            shutil.copytree(rendered_dir, keep_dir)
        return {
            "mode": "high-fidelity",
            "input": str(input_path),
            "output": str(output_path),
            "slides": len(analysis),
            "analysis": [a.__dict__ for a in analysis],
            "chrome": chrome,
            "node": node,
        }
    finally:
        if keep_images:
            tmp_root_obj.cleanup()
        else:
            tmp_root_obj.cleanup()


def main() -> None:
    import argparse

    ap = argparse.ArgumentParser(description="高保真 HTML slide deck → 南京模板 PPTX")
    ap.add_argument("--input", required=True, help="输入 HTML 文件")
    ap.add_argument("--output", required=True, help="输出 PPTX 文件")
    ap.add_argument("--chrome", default=None, help="Chrome/Edge/Chromium 可执行文件路径")
    ap.add_argument("--node", default=None, help="Node.js 可执行文件路径")
    ap.add_argument("--keep-images", action="store_true", help="保留渲染后的每页 PNG")
    args = ap.parse_args()

    report = render_html_to_pptx(
        args.input,
        args.output,
        chrome_path=args.chrome,
        node_path=args.node,
        keep_images=args.keep_images,
    )
    print("[OK] 已生成高保真南京模板 PPT:", report["output"])
    print("[INFO] 页数:", report["slides"])
    for item in report["analysis"]:
        print(
            f"  Slide {item['index']:>2}: {item['slide_type']:<7} | "
            f"density={item['density']:<6} | font_scale={item['font_scale']:.2f} | "
            f"text={item['text_length']:<4} | components={item['component_count']:<2} | {item['title'][:40]}"
        )


if __name__ == "__main__":
    main()
