#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
html_to_pptx.py — HTML 内容转腾讯云架构师南京城市沙龙模板 PPT (v2)

功能：
  解析 HTML 文件结构（标题/段落/列表/表格/代码块/图片/卡片/双栏/流程图/
  概率图/信息卡/策略对比/GPT大字/自定义flex布局），
  直接生成符合腾讯云架构师南京城市沙龙模板规范的 PPT。
  字体/配色/背景/Logo 均直接合规，无需二次适配。

字体大小标准（锁死规范，基于 20 英寸宽幅面 v3）：
  - 页面主标题: 40pt TencentSans W7 (章节页 48pt)
  - 编号(section-num): 18pt TencentSans W7
  - 内容副标题: 22pt TencentSans W3
  - 卡片标题: 26pt TencentSans W7
  - 列标题(双栏): 28pt TencentSans W7
  - 正文/卡片正文: 20pt TencentSans W3
  - 列表项: 20pt TencentSans W3
  - 流程图步骤: 22pt TencentSans W7
  - 流程图箭头: 28pt 南京主蓝
  - 高亮框: 20pt TencentSans W3
  - 策略对比 emoji: 42pt / label: 28pt W7 / desc: 20pt W3
  - 标签/脚注: 14-16pt TencentSans W7
  - 代码块: 14-16pt JetBrains Mono

支持的 HTML 元素映射：
  - H1 → 章节扇页（南京章节背景）
  - H2 → 内容页标题
  - H3/H4 → 内容页副标题
  - p → 正文段落
  - ul/ol/li → 列表（支持嵌套）
  - table → 表格页
  - pre/code → 代码块页
  - img → 图片页
  - blockquote → 引用框
  - .card-row / .card → 卡片布局
  - .two-col / .col → 双栏布局
  - .flow-row / .flow-step → 流程图
  - .highlight-box → 高亮框
  - .prob-row / .prob-bar → 概率条形图
  - .info-card → 信息卡片
  - .gpt-letters → 大字展示
  - .policy-grid / .policy-card → 策略对比
  - .summary-grid / .summary-card → 总结网格
  - inline-style flex containers → 自动识别并渲染

使用：
  python scripts/html_to_pptx.py --input article.html --output output.pptx
      [--title "演示标题"] [--subtitle "副标题"] [--author "作者"]

依赖：
  pip install python-pptx Pillow beautifulsoup4 lxml
"""

import argparse
import math
import os
import re
import sys
import urllib.request
import tempfile
from pathlib import Path
from typing import List, Tuple, Optional, Dict, Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    print("[ERROR] 缺少依赖。请运行: pip install python-pptx Pillow beautifulsoup4 lxml", file=sys.stderr)
    sys.exit(1)

try:
    from bs4 import BeautifulSoup, Tag, NavigableString
except ImportError:
    print("[ERROR] 缺少 beautifulsoup4。请运行: pip install beautifulsoup4", file=sys.stderr)
    sys.exit(1)

# ---------- 路径配置 ----------
SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent
ASSETS = SKILL_DIR / "assets"

BG_COVER = ASSETS / "backgrounds" / "bg-cover.jpeg"
BG_SECTION = ASSETS / "backgrounds" / "bg-section.jpeg"
BG_CONTENT = ASSETS / "backgrounds" / "bg-content.jpeg"
BG_END = ASSETS / "backgrounds" / "bg-end.jpeg"
LOGO_MAIN = ASSETS / "logos" / "logo-main.png"  # 南京模板默认不强制插入 Logo，保留变量兼容旧流程

# ---------- 品牌规范常量 ----------
# 页面尺寸（模板实测）
SLIDE_WIDTH_EMU = 18288000   # 20.00 inch
SLIDE_HEIGHT_EMU = 10287000  # 11.25 inch

# 字体
FONT_TITLE = "TencentSans W7"
FONT_BODY = "TencentSans W3"
FONT_CODE = "JetBrains Mono"

# 南京限定色板（只使用以下 7 色）
COLOR_BRAND_RED = RGBColor(0x32, 0x72, 0xDC)      # 兼容旧变量名：南京主蓝 #3272DC
COLOR_PRIMARY_BLUE = COLOR_BRAND_RED
COLOR_BLACK = RGBColor(0x08, 0x19, 0x4B)          # 南京深蓝 #08194B
COLOR_DARK_TEXT = COLOR_BLACK
COLOR_MID_GRAY = RGBColor(0x44, 0x47, 0x4F)       # #44474F
COLOR_LIGHT_GRAY = RGBColor(0x8B, 0x8C, 0x8C)     # #8B8C8C
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # #FFFFFF
COLOR_CREAM = COLOR_WHITE                         # 卡片底统一纯白，避免旧米白残留
COLOR_WARM_APRICOT = RGBColor(0x00, 0xC8, 0xD8)   # 兼容旧变量名：青蓝强调 #00C8D8
COLOR_CYAN = COLOR_WARM_APRICOT
COLOR_BRIGHT_BLUE = RGBColor(0x01, 0xA4, 0xFF)    # #01A4FF
COLOR_LINE_GRAY = COLOR_LIGHT_GRAY

# Logo 位置
LOGO_COVER_X = Inches(0.65)
LOGO_COVER_Y = Inches(0.61)
LOGO_CONTENT_X = Inches(14.60)
LOGO_CONTENT_Y = Inches(0.49)
LOGO_W = Inches(4.88)
LOGO_H = Inches(0.53)

# 内容区域
CONTENT_LEFT = Inches(1.0)
CONTENT_TOP = Inches(1.5)
CONTENT_WIDTH = Inches(18.0)
CONTENT_HEIGHT = Inches(9.0)

# 每页最大内容行数（用于智能分页）
MAX_LINES_PER_SLIDE = 12
MAX_LIST_ITEMS_PER_SLIDE = 10
MAX_TABLE_ROWS_PER_SLIDE = 8


# ---------- HTML 解析 ----------

class ContentBlock:
    """解析后的内容块"""
    def __init__(self, block_type: str, **kwargs):
        self.type = block_type  # h1, h2, h3, paragraph, list, table, code, image, blockquote
        self.text = kwargs.get('text', '')
        self.items = kwargs.get('items', [])  # list items
        self.rows = kwargs.get('rows', [])    # table rows
        self.headers = kwargs.get('headers', [])  # table headers
        self.language = kwargs.get('language', '')  # code language
        self.src = kwargs.get('src', '')      # image source
        self.alt = kwargs.get('alt', '')      # image alt text
        self.level = kwargs.get('level', 0)   # heading level / list nesting
        self.ordered = kwargs.get('ordered', False)  # ordered list


def is_slide_deck_html(html_content: str) -> bool:
    """检测 HTML 是否是幻灯片演示型（如 reveal.js / impress.js / 自定义 slide deck）。

    检测标准：存在 3 个以上具有 class="slide" 或 data-slide 属性的 div。
    """
    soup = BeautifulSoup(html_content, 'lxml')
    slides = soup.find_all('div', class_=re.compile(r'\bslide\b'))
    if len(slides) >= 3:
        return True
    slides_by_attr = soup.find_all('div', attrs={'data-slide': True})
    if len(slides_by_attr) >= 3:
        return True
    # 也检查 section.slide（reveal.js 风格）
    slides_section = soup.find_all('section', class_=re.compile(r'\bslide\b'))
    if len(slides_section) >= 3:
        return True
    return False


class SlideContent:
    """一个 HTML 幻灯片页面的结构化内容"""
    def __init__(self):
        self.slide_type = 'content'  # cover, section, content, end
        self.section_num = ''
        self.title = ''
        self.subtitle = ''
        self.author = ''
        self.tags = []
        self.cards = []         # [{title, body, tag, style, icon}]
        self.columns = []       # [{title, text, items, highlight, cards}]
        self.flow_steps = []    # [{text, style}]
        self.highlights = []    # [text]
        self.paragraphs = []    # [text]
        self.list_items = []    # [text]
        self.policy_cards = []  # [{who, label, desc}]
        self.summary_cards = [] # [{label, title, desc, analogy}]
        self.prob_charts = []   # [{title, items: [{label, pct, color}]}]
        self.info_cards = []    # [{title, body, style}]
        self.gpt_letters = []   # [{letter, color, desc}]
        self.inline_flex_blocks = []  # [{items: [{icon, title, desc}]}]
        self.raw_text_blocks = []  # 无法分类的文字块


def _extract_card_data(card) -> dict:
    """从一个 .card 元素中提取结构化数据
    
    支持多种卡片结构：
    1. 标准: .card-icon + .card-title + .card-body
    2. 内联样式: font-weight:800 div 做标题 + .col-text 做正文
    3. 纯文本: 无特殊class时提取全部文本
    """
    card_data = {}
    ct = card.find(class_='card-title')
    cb = card.find(class_='card-body')
    ci = card.find(class_='card-icon')
    ctag = card.find(class_='card-tag')
    card_data['title'] = ct.get_text(strip=True) if ct else ''
    card_data['body'] = _clean_html_text(cb) if cb else ''
    card_data['icon'] = ci.get_text(strip=True) if ci else ''
    card_data['tag'] = ctag.get_text(strip=True) if ctag else ''
    
    # 如果没有 card-body，尝试从 col-text 提取（Slide 10/13 的卡片结构）
    if not card_data['body']:
        col_text = card.find(class_='col-text')
        if col_text:
            card_data['body'] = _clean_html_text(col_text)
    
    # 如果仍无 body，尝试 col-list
    if not card_data['body']:
        col_list = card.find(class_='col-list')
        if col_list:
            items = [li.get_text(strip=True) for li in col_list.find_all('li')]
            card_data['body'] = '\n'.join(f'• {item}' for item in items)
    
    # 如果仍无 body 也无 title，提取全部文本
    if not card_data['body'] and not card_data['title']:
        card_data['body'] = _clean_html_text(card)
    
    # 如果有内联标题（font-bold / font-weight:800 div 做标题）但没有 card-title
    if not card_data['title']:
        for div in card.find_all('div', recursive=False):
            style = div.get('style', '')
            div_classes = ' '.join(div.get('class', []))
            # 跳过 col-text（那是正文）
            if 'col-text' in div_classes:
                continue
            if ('font-bold' in style or 'font-weight:800' in style or 
                'font-weight: 800' in style or 'font-weight:700' in style):
                card_data['title'] = div.get_text(strip=True)
                break
    
    # 如果 body 里和 title 完全重复（因为全文提取包含了标题），去重
    if card_data['title'] and card_data['body'] and card_data['body'].startswith(card_data['title']):
        remaining = card_data['body'][len(card_data['title']):].strip()
        if remaining:
            card_data['body'] = remaining
    
    # 检测卡片风格（class优先，inline-style兜底）
    card_classes = ' '.join(card.get('class', []))
    card_inline_style = card.get('style', '')
    if 'accent' in card_classes:
        card_data['style'] = 'accent'
    elif 'teal' in card_classes:
        card_data['style'] = 'teal'
    elif 'purple' in card_classes:
        card_data['style'] = 'purple'
    elif 'blue' in card_classes:
        card_data['style'] = 'blue'
    elif 'gold' in card_classes:
        card_data['style'] = 'gold'
    elif 'red' in card_inline_style or '--red' in card_inline_style or '--accent' in card_inline_style:
        card_data['style'] = 'accent'
    elif '--blue' in card_inline_style:
        card_data['style'] = 'blue'
    else:
        card_data['style'] = 'default'
    return card_data


def parse_html_slides(html_content: str) -> List[SlideContent]:
    """解析幻灯片型 HTML 为 SlideContent 列表"""
    soup = BeautifulSoup(html_content, 'lxml')

    # 找所有 slide 元素
    slides = soup.find_all('div', class_=re.compile(r'\bslide\b'))
    if not slides:
        slides = soup.find_all('div', attrs={'data-slide': True})
    if not slides:
        slides = soup.find_all('section', class_=re.compile(r'\bslide\b'))

    results = []
    for slide_elem in slides:
        sc = SlideContent()

        # ---------- 判断页面类型 ----------
        classes = ' '.join(slide_elem.get('class', []))
        if 'title-page' in classes or 'slide-title' in classes:
            sc.slide_type = 'cover'
        elif 'slide-center' in classes:
            # 可能是章节页或结尾页，看内容判断
            text_content = slide_elem.get_text(strip=True)
            if '谢谢' in text_content or 'Thanks' in text_content.lower() or 'thank' in text_content.lower():
                sc.slide_type = 'end'
            else:
                sc.slide_type = 'section'

        # ---------- 提取编号 ----------
        section_num = slide_elem.find(class_='section-num')
        if section_num:
            sc.section_num = section_num.get_text(strip=True)

        # ---------- 提取标题 ----------
        title_main = slide_elem.find(class_='title-main')
        section_heading = slide_elem.find(class_='section-heading')
        if title_main:
            sc.title = title_main.get_text(strip=True)
            sc.slide_type = 'cover'
        elif section_heading:
            sc.title = section_heading.get_text(strip=True)

        # ---------- 提取副标题 ----------
        sub = slide_elem.find(class_='section-sub') or slide_elem.find(class_='title-sub')
        if sub:
            sc.subtitle = sub.get_text(strip=True)

        # ---------- 提取作者 ----------
        # 封面页特有
        if sc.slide_type == 'cover':
            # 在封面中找作者信息（通常在 title 下面的 div 中）
            for div in slide_elem.find_all('div', recursive=True):
                text = div.get_text(strip=True)
                if '|' in text and len(text) < 100:
                    sc.author = text
                    break

        # ---------- 提取标签 ----------
        tags_container = slide_elem.find(class_='title-tags')
        if tags_container:
            for tag in tags_container.find_all(class_='title-tag'):
                sc.tags.append(tag.get_text(strip=True))
        else:
            # 也检查没有 title-tags 容器但散落的 title-tag 元素
            for tag in slide_elem.find_all(class_='title-tag'):
                tag_text = tag.get_text(strip=True)
                if tag_text and tag_text not in sc.tags:
                    sc.tags.append(tag_text)

        # ---------- 提取卡片（card-row 内和独立 card） ----------
        card_rows = slide_elem.find_all(class_='card-row')
        for card_row in card_rows:
            for card in card_row.find_all(class_='card'):
                card_data = _extract_card_data(card)
                sc.cards.append(card_data)
        
        # 独立的 .card 不在 card-row 内（如 Slide 7/11/12 的自由布局）
        if not card_rows:
            standalone_cards = slide_elem.find_all(class_='card', recursive=True)
            # 排除在 two-col 内的 card（那些归双栏处理）
            two_col_elem = slide_elem.find(class_='two-col')
            for card in standalone_cards:
                # 跳过在 two-col 内的
                if two_col_elem and card.find_parent(class_='two-col'):
                    continue
                card_data = _extract_card_data(card)
                sc.cards.append(card_data)

        # ---------- 提取双栏 ----------
        two_col = slide_elem.find(class_='two-col')
        if two_col:
            for col in two_col.find_all(class_='col', recursive=False):
                col_data = {}
                ct = col.find(class_='col-title')
                col_data['title'] = ct.get_text(strip=True) if ct else ''
                
                # 内部卡片 (.card) — 优先提取
                col_cards = col.find_all(class_='card', recursive=True)
                col_data['cards'] = []
                for cc in col_cards:
                    col_data['cards'].append(_extract_card_data(cc))
                
                # 获取列表（可能多个 col-list，排除在 card 内的）
                all_items = []
                for col_list in col.find_all(class_='col-list'):
                    # 只取不在 card 内的 col-list
                    if col_list.find_parent(class_='card'):
                        continue
                    for li in col_list.find_all('li'):
                        all_items.append(li.get_text(strip=True))
                col_data['items'] = all_items
                
                # 获取正文（排除在 card 内的 col-text）
                col_data['text'] = ''
                for ct_elem in col.find_all(class_='col-text'):
                    if ct_elem.find_parent(class_='card'):
                        continue  # card 内的 col-text 已由 _extract_card_data 处理
                    col_data['text'] = _clean_html_text(ct_elem)
                    break
                
                # 高亮框（排除在 card 内的，多个）
                highlight_texts = []
                for hb in col.find_all(class_='highlight-box'):
                    if hb.find_parent(class_='card'):
                        continue
                    ht = _clean_html_text(hb)
                    if ht:
                        highlight_texts.append(ht)
                col_data['highlight'] = '\n'.join(highlight_texts) if highlight_texts else ''
                
                # 内部的独立div（非 card / col-text / col-list / highlight-box）
                inner_boxes = col.find_all('div', recursive=False)
                inner_texts = []
                for ib in inner_boxes:
                    ib_classes = ' '.join(ib.get('class', []))
                    if any(k in ib_classes for k in ['col-title', 'col-text', 'col-list',
                                                      'highlight-box', 'big-num', 'card']):
                        continue
                    ib_text = _clean_html_text(ib)
                    if ib_text and len(ib_text) > 10:
                        inner_texts.append(ib_text)
                if inner_texts and not col_data['text'] and not col_data['items'] and not col_data['cards']:
                    col_data['text'] = '\n'.join(inner_texts)
                sc.columns.append(col_data)

        # ---------- 提取流程图 ----------
        flow_rows = slide_elem.find_all(class_='flow-row')
        for flow_row in flow_rows:
            for step in flow_row.find_all(class_='flow-step'):
                step_data = {}
                step_data['text'] = step.get_text(strip=True)
                step_classes = ' '.join(step.get('class', []))
                if 'accent' in step_classes:
                    step_data['style'] = 'accent'
                elif 'teal' in step_classes:
                    step_data['style'] = 'teal'
                elif 'purple' in step_classes:
                    step_data['style'] = 'purple'
                else:
                    step_data['style'] = 'default'
                sc.flow_steps.append(step_data)

        # ---------- 提取高亮框 ----------
        # 只取不在 two-col 或 card 内的高亮框
        for hb in slide_elem.find_all(class_='highlight-box', recursive=True):
            # 检查是否已被双栏收走
            parent_col = hb.find_parent(class_='col')
            if parent_col and two_col and parent_col in two_col.find_all(class_='col'):
                continue
            sc.highlights.append(_clean_html_text(hb))

        # ---------- 提取策略对比 ----------
        policy_grid = slide_elem.find(class_='policy-grid')
        if policy_grid:
            for pc in policy_grid.find_all(class_='policy-card'):
                pdata = {}
                who = pc.find(class_='who')
                label = pc.find(class_='label')
                pdata['who'] = who.get_text(strip=True) if who else ''
                pdata['label'] = label.get_text(strip=True) if label else ''
                # 提取大 emoji（通常是 font-size:42px 的 div）
                pdata['emoji'] = ''
                for d in pc.find_all('div', recursive=False):
                    style = d.get('style', '').replace(' ', '')
                    if 'font-size:42px' in style or 'font-size:48px' in style:
                        pdata['emoji'] = d.get_text(strip=True)
                        break
                # 描述（排除 who / label / emoji）
                desc_parts = []
                for d in pc.find_all('div', recursive=False):
                    dc = ' '.join(d.get('class', []))
                    if 'who' in dc or 'label' in dc:
                        continue
                    style = d.get('style', '').replace(' ', '')
                    # 跳过 emoji div
                    if 'font-size:42px' in style or 'font-size:48px' in style:
                        continue
                    t = d.get_text(strip=True)
                    if t and len(t) > 3:
                        desc_parts.append(t)
                pdata['desc'] = '\n'.join(desc_parts)
                sc.policy_cards.append(pdata)

        # ---------- 提取总结网格 ----------
        summary_grid = slide_elem.find(class_='summary-grid')
        if summary_grid:
            for sc_card in summary_grid.find_all(class_='summary-card'):
                sdata = {}
                ml = sc_card.find(class_='method-label')
                mt = sc_card.find(class_='method-title')
                md = sc_card.find(class_='method-desc')
                ma = sc_card.find(class_='method-analogy')
                sdata['label'] = ml.get_text(strip=True) if ml else ''
                sdata['title'] = mt.get_text(strip=True) if mt else ''
                sdata['desc'] = _clean_html_text(md) if md else ''
                sdata['analogy'] = ma.get_text(strip=True) if ma else ''
                sc.summary_cards.append(sdata)

        # ---------- 提取概率条形图 ----------
        prob_rows = slide_elem.find_all(class_='prob-row')
        if prob_rows:
            chart_data = {'title': '', 'items': []}
            # 找图表标题（通常在 prob-row 前面的 info-title 或 div 中）
            prob_parent = prob_rows[0].parent
            if prob_parent:
                title_elem = prob_parent.find(class_='info-title')
                if title_elem:
                    chart_data['title'] = title_elem.get_text(strip=True)
                # 找描述
                desc_elem = prob_parent.find(string=re.compile(r'输入'))
                if desc_elem:
                    chart_data['title'] += f' ({desc_elem.strip()})'
            for pr in prob_rows:
                label_elem = pr.find(class_='prob-label')
                pct_elem = pr.find(class_='prob-pct')
                label = label_elem.get_text(strip=True) if label_elem else ''
                pct_text = pct_elem.get_text(strip=True) if pct_elem else '0%'
                chart_data['items'].append({
                    'label': label,
                    'pct': pct_text,
                })
            sc.prob_charts.append(chart_data)

        # ---------- 提取信息卡片（info-card）----------
        info_cards = slide_elem.find_all(class_='info-card')
        # 排除在 prob-chart 内的父容器
        for ic in info_cards:
            # 跳过包含 prob-row 的容器（那是概率图的容器）
            if ic.find(class_='prob-row'):
                continue
            idata = {}
            it = ic.find(class_='info-title')
            ib = ic.find(class_='info-body')
            idata['title'] = it.get_text(strip=True) if it else ''
            idata['body'] = _clean_html_text(ib) if ib else _clean_html_text(ic)
            ic_classes = ' '.join(ic.get('class', []))
            if 'accent' in ic_classes:
                idata['style'] = 'accent'
            elif 'blue' in ic_classes:
                idata['style'] = 'blue'
            else:
                idata['style'] = 'default'
            sc.info_cards.append(idata)

        # ---------- 提取 GPT 大字（gpt-letters）----------
        gpt_elem = slide_elem.find(class_='gpt-letters')
        if gpt_elem:
            for span in gpt_elem.find_all('span'):
                letter = span.get_text(strip=True)
                if letter:
                    sc.gpt_letters.append({'letter': letter})
            # 提取大字下方的解释块
            gpt_parent = gpt_elem.parent
            if gpt_parent:
                flex_divs = gpt_parent.find_all('div', recursive=False)
                for fd in flex_divs:
                    style = fd.get('style', '')
                    if 'display:flex' in style.replace(' ', '') or 'display: flex' in style:
                        # 这是解释块的容器
                        items = []
                        for child_div in fd.find_all('div', recursive=False):
                            item = {}
                            # 找图标（大 emoji）
                            all_divs = child_div.find_all('div', recursive=False)
                            for d in all_divs:
                                text = d.get_text(strip=True)
                                d_style = d.get('style', '')
                                if 'font-size:52px' in d_style.replace(' ', '') or 'font-size: 52px' in d_style:
                                    item['icon'] = text
                                elif 'font-bold' in d_style or 'font-weight:800' in d_style.replace(' ', ''):
                                    item['title'] = text
                                elif 'font-size:19px' in d_style.replace(' ', '') or 'font-size: 19px' in d_style:
                                    item['desc'] = text
                            if item.get('title'):
                                items.append(item)
                        if items:
                            sc.inline_flex_blocks.append({'items': items})

        # ---------- 提取剩余文本内容 ----------
        # 只有当所有结构化提取都为空时才收集剩余文本
        has_structured = (sc.cards or sc.columns or sc.summary_cards or sc.policy_cards
                         or sc.prob_charts or sc.info_cards or sc.gpt_letters
                         or sc.inline_flex_blocks)
        if not has_structured:
            # 排除已提取的部分
            excluded_classes = {'section-num', 'section-heading', 'section-sub',
                              'title-main', 'title-sub', 'title-badge', 'title-tags',
                              'glow-top-right', 'glow-bottom-left', 'nav',
                              'flow-row', 'highlight-box', 'card-row', 'two-col',
                              'policy-grid', 'summary-grid', 'page-header',
                              'prob-row', 'info-card', 'gpt-letters'}
            for child in slide_elem.find_all('div', recursive=False):
                child_classes = set(child.get('class', []))
                if child_classes & excluded_classes:
                    continue
                if child_classes == {'glow-top-right'} or child_classes == {'glow-bottom-left'}:
                    continue
                # 跳过已经是 page-header 的
                if 'page-header' in child_classes:
                    continue
                text = _clean_html_text(child)
                if text and len(text) > 5:
                    sc.raw_text_blocks.append(text)

        results.append(sc)

    return results


def _clean_html_text(elem) -> str:
    """从元素中提取干净的文本，保留基本换行"""
    if elem is None:
        return ''
    # 用 <br> 标记换行
    for br in elem.find_all('br'):
        br.replace_with('\n')
    text = elem.get_text(separator='\n')
    # 清理多余空行
    lines = [line.strip() for line in text.split('\n')]
    lines = [line for line in lines if line]
    return '\n'.join(lines)


def parse_html(html_content: str) -> List[ContentBlock]:
    """解析 HTML 内容为结构化内容块列表（文章型 HTML）"""
    soup = BeautifulSoup(html_content, 'lxml')
    
    # 尝试找到主要内容区域
    main = soup.find('main') or soup.find('article') or soup.find('body') or soup
    
    blocks = []
    _parse_elements(main, blocks)
    return blocks


def _parse_elements(parent, blocks: List[ContentBlock]):
    """递归解析 DOM 元素"""
    for elem in parent.children:
        if isinstance(elem, NavigableString):
            text = str(elem).strip()
            if text and text != '\n':
                blocks.append(ContentBlock('paragraph', text=text))
            continue
        
        if not isinstance(elem, Tag):
            continue
        
        tag = elem.name.lower()
        
        # 标题
        if tag in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6'):
            level = int(tag[1])
            text = elem.get_text(strip=True)
            if text:
                blocks.append(ContentBlock(f'h{level}', text=text, level=level))
        
        # 段落
        elif tag == 'p':
            text = elem.get_text(strip=True)
            if text:
                # 检查是否包含图片
                img = elem.find('img')
                if img:
                    src = img.get('src', '')
                    alt = img.get('alt', '')
                    blocks.append(ContentBlock('image', src=src, alt=alt))
                else:
                    blocks.append(ContentBlock('paragraph', text=text))
        
        # 列表
        elif tag in ('ul', 'ol'):
            items = _parse_list_items(elem)
            if items:
                blocks.append(ContentBlock('list', items=items, ordered=(tag == 'ol')))
        
        # 表格
        elif tag == 'table':
            headers, rows = _parse_table(elem)
            if headers or rows:
                blocks.append(ContentBlock('table', headers=headers, rows=rows))
        
        # 代码块
        elif tag == 'pre':
            code_elem = elem.find('code')
            if code_elem:
                text = code_elem.get_text()
                lang_class = code_elem.get('class', [])
                language = ''
                if lang_class:
                    for cls in lang_class:
                        if cls.startswith('language-') or cls.startswith('lang-'):
                            language = cls.split('-', 1)[1]
                            break
                blocks.append(ContentBlock('code', text=text, language=language))
            else:
                text = elem.get_text()
                blocks.append(ContentBlock('code', text=text))
        
        # 引用
        elif tag == 'blockquote':
            text = elem.get_text(strip=True)
            if text:
                blocks.append(ContentBlock('blockquote', text=text))
        
        # 图片
        elif tag == 'img':
            src = elem.get('src', '')
            alt = elem.get('alt', '')
            if src:
                blocks.append(ContentBlock('image', src=src, alt=alt))
        
        # div/section/article 等容器：递归处理
        elif tag in ('div', 'section', 'article', 'header', 'footer', 'main', 'aside', 'nav', 'span', 'strong', 'em', 'b', 'i'):
            _parse_elements(elem, blocks)


def _parse_list_items(list_elem) -> List[str]:
    """解析列表项"""
    items = []
    for li in list_elem.find_all('li', recursive=False):
        text = ''
        # 获取直接文本内容（不含嵌套列表）
        for child in li.children:
            if isinstance(child, NavigableString):
                text += str(child)
            elif isinstance(child, Tag) and child.name not in ('ul', 'ol'):
                text += child.get_text()
        text = text.strip()
        if text:
            items.append(text)
        # 嵌套列表展平（前加缩进标记）
        nested = li.find(['ul', 'ol'])
        if nested:
            for sub_li in nested.find_all('li', recursive=False):
                sub_text = sub_li.get_text(strip=True)
                if sub_text:
                    items.append('  • ' + sub_text)
    return items


def _parse_table(table_elem) -> Tuple[List[str], List[List[str]]]:
    """解析表格"""
    headers = []
    rows = []
    
    # 表头
    thead = table_elem.find('thead')
    if thead:
        tr = thead.find('tr')
        if tr:
            headers = [th.get_text(strip=True) for th in tr.find_all(['th', 'td'])]
    
    # 表体
    tbody = table_elem.find('tbody') or table_elem
    for tr in tbody.find_all('tr'):
        cells = [td.get_text(strip=True) for td in tr.find_all(['td', 'th'])]
        if cells:
            # 如果没有 thead 但第一行是 th，视为表头
            if not headers and tr.find('th') and len(rows) == 0:
                headers = cells
            else:
                rows.append(cells)
    
    return headers, rows


# ---------- PPT 生成 ----------

class PptxGenerator:
    """PPT 生成器"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH_EMU
        self.prs.slide_height = SLIDE_HEIGHT_EMU
        self.slide_count = 0
        self._blank_layout = self.prs.slide_layouts[6]  # 空白版式
    
    def generate_from_slides(self, slide_contents: List[SlideContent]) -> None:
        """从幻灯片型 HTML 解析结果生成 PPT（保留原始 slide 结构）"""
        for sc in slide_contents:
            if sc.slide_type == 'cover':
                self._add_cover_slide(sc.title, sc.subtitle, sc.author)
                # 添加标签
                if sc.tags:
                    slide = self.prs.slides[-1]
                    self._add_tags_row(slide, sc.tags, Inches(2.0), Inches(8.5))
            elif sc.slide_type == 'end':
                self._add_end_slide_custom(sc.title, sc.tags)
            else:
                # 内容页/章节页
                self._add_structured_slide(sc)

    def _add_end_slide_custom(self, title: str = '', tags: list = None):
        """自定义结尾页"""
        slide = self._new_slide('end')
        
        display_title = title if title else "THANKS"
        left = Inches(5.0)
        top = Inches(3.5)
        width = Inches(10.0)
        height = Inches(2.5)
        tf = self._add_textbox(slide, left, top, width, height)
        self._add_text_run(tf, display_title, font_name=FONT_TITLE, font_size=Pt(60),
                         color=COLOR_BRAND_RED, bold=True, alignment=PP_ALIGN.CENTER)
        
        if tags:
            self._add_tags_row(slide, tags, Inches(4.0), Inches(7.0))

    def _add_tags_row(self, slide, tags: list, left, top):
        """添加标签行"""
        tag_text = '    '.join(f'[{t}]' for t in tags)
        width = Inches(14.0)
        height = Inches(0.8)
        tf = self._add_textbox(slide, left, top, width, height)
        self._add_text_run(tf, tag_text, font_name=FONT_BODY, font_size=Pt(18),
                         color=COLOR_MID_GRAY, alignment=PP_ALIGN.LEFT)

    def _add_structured_slide(self, sc: SlideContent):
        """生成结构化内容幻灯片（保留卡片/双栏/流程图/概率图/信息卡等布局）"""
        # 判断是否应该作为章节页
        has_content = (sc.cards or sc.columns or sc.flow_steps or
                      sc.highlights or sc.summary_cards or sc.policy_cards or
                      sc.raw_text_blocks or sc.prob_charts or sc.info_cards or
                      sc.gpt_letters or sc.inline_flex_blocks)
        is_section = (sc.slide_type == 'section') or (
            not has_content and sc.title and not sc.subtitle
        )
        
        if is_section and not sc.cards and not sc.columns:
            slide = self._new_slide('section')
            page_type = 'section'
        else:
            slide = self._new_slide('content')
            page_type = 'content'

        y_cursor = Inches(1.3)

        # ---------- 编号 + 标题 ----------
        if sc.section_num:
            left = Inches(1.0)
            width = Inches(16.0)
            height = Inches(0.5)
            tf = self._add_textbox(slide, left, y_cursor, width, height)
            self._add_text_run(tf, sc.section_num, font_name=FONT_TITLE, font_size=Pt(18),
                             color=COLOR_BRAND_RED, bold=True, alignment=PP_ALIGN.LEFT)
            y_cursor += Inches(0.5)

        if sc.title:
            left = Inches(1.0)
            width = Inches(17.0)
            height = Inches(1.3)
            tf = self._add_textbox(slide, left, y_cursor, width, height)
            size = Pt(48) if page_type == 'section' else Pt(40)
            self._add_text_run(tf, sc.title, font_name=FONT_TITLE, font_size=size,
                             color=COLOR_DARK_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
            y_cursor += Inches(1.1) if page_type == 'section' else Inches(1.0)

        if sc.subtitle:
            left = Inches(1.0)
            width = Inches(16.0)
            height = Inches(0.7)
            tf = self._add_textbox(slide, left, y_cursor, width, height)
            self._add_text_run(tf, sc.subtitle, font_name=FONT_BODY, font_size=Pt(22),
                             color=COLOR_MID_GRAY, alignment=PP_ALIGN.LEFT)
            y_cursor += Inches(0.7)

        # ---------- 高亮框（置顶，如 Slide 12 顶部的策略定义框） ----------
        if sc.highlights:
            for hl_text in sc.highlights:
                if y_cursor > Inches(9.5):
                    break
                y_cursor = self._render_highlight_box(slide, hl_text, y_cursor)

        # ---------- GPT 大字展示 ----------
        if sc.gpt_letters:
            y_cursor = self._render_gpt_letters(slide, sc.gpt_letters, y_cursor)

        # ---------- 内联 Flex 块（大字下方的解释） ----------
        if sc.inline_flex_blocks:
            for ifb in sc.inline_flex_blocks:
                y_cursor = self._render_inline_flex(slide, ifb['items'], y_cursor)

        # ---------- 卡片布局（不互斥，允许和 policy_cards 共存） ----------
        if sc.cards:
            y_cursor = self._render_cards(slide, sc.cards, y_cursor)

        # ---------- 双栏布局 ----------
        if sc.columns:
            y_cursor = self._render_columns(slide, sc.columns, y_cursor)

        # ---------- 策略对比（不再 elif，允许和 cards 共存） ----------
        if sc.policy_cards:
            y_cursor = self._render_policy_grid(slide, sc.policy_cards, y_cursor)

        # ---------- 总结网格 ----------
        if sc.summary_cards:
            y_cursor = self._render_summary_grid(slide, sc.summary_cards, y_cursor)

        # ---------- 剩余文本内容（只有真的没有其他结构时） ----------
        has_main_content = (sc.cards or sc.columns or sc.policy_cards or
                           sc.summary_cards or sc.gpt_letters or sc.inline_flex_blocks)
        if sc.raw_text_blocks and not has_main_content:
            for text_block in sc.raw_text_blocks:
                if y_cursor > Inches(9.0):
                    break
                lines = text_block.count('\n') + 1
                height = Inches(min(max(0.6, lines * 0.4), 4.0))
                left = Inches(1.0)
                width = Inches(17.0)
                tf = self._add_textbox(slide, left, y_cursor, width, height)
                self._add_text_run(tf, text_block, font_name=FONT_BODY, font_size=Pt(20),
                                 color=COLOR_DARK_TEXT, alignment=PP_ALIGN.LEFT)
                y_cursor += height + Inches(0.2)

        # ---------- 流程图（在所有主体内容下方） ----------
        if sc.flow_steps:
            y_cursor = self._render_flow(slide, sc.flow_steps, y_cursor)

        # ---------- 概率条形图 ----------
        if sc.prob_charts:
            for chart in sc.prob_charts:
                y_cursor = self._render_prob_chart(slide, chart, y_cursor)

        # ---------- 信息卡片 ----------
        if sc.info_cards:
            y_cursor = self._render_info_cards(slide, sc.info_cards, y_cursor)

    def _render_cards(self, slide, cards: list, y_start) -> int:
        """渲染卡片行（字体规范：标题26pt W7，正文20pt W3，标签16pt W7）"""
        n_cards = len(cards)
        if n_cards == 0:
            return y_start

        card_gap = Inches(0.3)
        total_width = Inches(17.5)
        card_left_start = Inches(1.0)
        card_width = int((total_width - card_gap * (n_cards - 1)) / n_cards)
        
        # 根据内容量动态计算卡片高度
        max_body_lines = max((len(c.get('body', '').split('\n')) for c in cards), default=3)
        card_height = Inches(min(max(3.5, max_body_lines * 0.45 + 2.0), 6.0))
        
        # 如果剩余空间不够，缩减高度
        available = Inches(9.5) - y_start
        if card_height > available:
            card_height = available

        for idx, card in enumerate(cards):
            left = card_left_start + idx * (card_width + card_gap)
            
            # 卡片背景框（根据风格设置颜色）
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y_start, card_width, card_height)
            shape.fill.solid()
            card_style = card.get('style', 'default')
            if card_style == 'accent':
                shape.fill.fore_color.rgb = COLOR_WHITE
                shape.line.color.rgb = COLOR_BRAND_RED
            elif card_style == 'blue':
                shape.fill.fore_color.rgb = COLOR_WHITE
                shape.line.color.rgb = COLOR_PRIMARY_BLUE
            else:
                shape.fill.fore_color.rgb = COLOR_CREAM
                shape.line.color.rgb = COLOR_LINE_GRAY
            shape.adjustments[0] = 0.04

            # 卡片内容文本框
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.35)
            tf.margin_right = Inches(0.35)
            tf.margin_top = Inches(0.3)
            tf.margin_bottom = Inches(0.3)
            tf.vertical_anchor = MSO_ANCHOR.TOP

            # 图标
            if card.get('icon'):
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = card['icon']
                run.font.size = Pt(36)

            # 标题 — 26pt TencentSans W7
            if card.get('title'):
                p = tf.add_paragraph()
                p.space_before = Pt(10)
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = card['title']
                run.font.name = FONT_TITLE
                run.font.size = Pt(26)
                run.font.bold = True
                # 标题颜色跟随卡片风格
                if card_style == 'accent':
                    run.font.color.rgb = COLOR_BRAND_RED
                elif card_style == 'blue':
                    run.font.color.rgb = COLOR_PRIMARY_BLUE
                else:
                    run.font.color.rgb = COLOR_DARK_TEXT
                self._set_font_east_asian(run, FONT_TITLE)

            # 正文 — 20pt TencentSans W3（卡片内正文）
            if card.get('body'):
                body_text = card['body']
                if len(body_text) > 400:
                    body_text = body_text[:400] + '...'
                # 按行分段渲染
                body_lines = body_text.split('\n')
                for line in body_lines[:14]:  # 最多14行
                    if not line.strip():
                        continue
                    p = tf.add_paragraph()
                    p.space_before = Pt(5)
                    p.alignment = PP_ALIGN.LEFT
                    run = p.add_run()
                    run.text = line
                    run.font.name = FONT_BODY
                    run.font.size = Pt(20)
                    run.font.color.rgb = COLOR_DARK_TEXT
                    self._set_font_east_asian(run, FONT_BODY)

            # 底部标签 — 16pt W7 南京主蓝
            if card.get('tag'):
                p = tf.add_paragraph()
                p.space_before = Pt(14)
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = f"[{card['tag']}]"
                run.font.name = FONT_TITLE
                run.font.size = Pt(16)
                run.font.color.rgb = COLOR_BRAND_RED
                run.font.bold = True
                self._set_font_east_asian(run, FONT_TITLE)

        return y_start + card_height + Inches(0.3)

    def _render_columns(self, slide, columns: list, y_start) -> int:
        """渲染双栏布局（字体规范：列标题28pt W7，列表项20pt W3，正文20pt W3）"""
        n_cols = len(columns)
        if n_cols == 0:
            return y_start

        col_gap = Inches(0.4)
        total_width = Inches(17.5)
        col_left_start = Inches(1.0)
        col_width = int((total_width - col_gap * (n_cols - 1)) / n_cols)
        col_height = Inches(5.8)
        
        available = Inches(9.2) - y_start
        if col_height > available:
            col_height = available

        for idx, col in enumerate(columns):
            left = col_left_start + idx * (col_width + col_gap)
            
            # 判断列背景色（根据内部 card 的风格）
            col_fill_color = COLOR_CREAM
            col_line_color = COLOR_LINE_GRAY
            if col.get('cards'):
                first_card_style = col['cards'][0].get('style', 'default')
                if first_card_style == 'accent':
                    col_fill_color = COLOR_WHITE
                    col_line_color = COLOR_BRAND_RED
                elif first_card_style == 'blue':
                    col_fill_color = COLOR_WHITE
                    col_line_color = COLOR_PRIMARY_BLUE
            
            # 列背景框
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y_start, col_width, col_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = col_fill_color
            shape.line.color.rgb = col_line_color
            shape.adjustments[0] = 0.04

            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.3)
            tf.margin_right = Inches(0.3)
            tf.margin_top = Inches(0.25)
            tf.margin_bottom = Inches(0.25)
            tf.vertical_anchor = MSO_ANCHOR.TOP

            # 列标题 — 28pt TencentSans W7
            if col.get('title'):
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = col['title']
                run.font.name = FONT_TITLE
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = COLOR_DARK_TEXT
                self._set_font_east_asian(run, FONT_TITLE)

            # 优先渲染内部卡片（如 Slide 10/13 的双栏中嵌套 card 结构）
            if col.get('cards'):
                for cc in col['cards'][:2]:
                    # 卡片标题 — 24pt W7
                    if cc.get('title'):
                        p = tf.add_paragraph()
                        p.space_before = Pt(12)
                        p.alignment = PP_ALIGN.LEFT
                        run = p.add_run()
                        run.text = cc['title']
                        run.font.name = FONT_TITLE
                        run.font.size = Pt(24)
                        run.font.bold = True
                        # 根据卡片风格设置标题颜色
                        if cc.get('style') == 'accent':
                            run.font.color.rgb = COLOR_BRAND_RED
                        elif cc.get('style') == 'blue':
                            run.font.color.rgb = COLOR_PRIMARY_BLUE
                        else:
                            run.font.color.rgb = COLOR_DARK_TEXT
                        self._set_font_east_asian(run, FONT_TITLE)
                    # 卡片正文 — 20pt W3
                    if cc.get('body'):
                        body = cc['body']
                        if len(body) > 350:
                            body = body[:350] + '...'
                        for line in body.split('\n')[:10]:
                            if not line.strip():
                                continue
                            p = tf.add_paragraph()
                            p.space_before = Pt(5)
                            p.alignment = PP_ALIGN.LEFT
                            run = p.add_run()
                            run.text = line
                            run.font.name = FONT_BODY
                            run.font.size = Pt(20)
                            run.font.color.rgb = COLOR_DARK_TEXT
                            self._set_font_east_asian(run, FONT_BODY)

            # 列表项 — 20pt TencentSans W3
            elif col.get('items'):
                for item in col['items'][:10]:  # 最多10项
                    p = tf.add_paragraph()
                    p.space_before = Pt(6)
                    p.alignment = PP_ALIGN.LEFT
                    run = p.add_run()
                    run.text = f"• {item}"
                    run.font.name = FONT_BODY
                    run.font.size = Pt(20)
                    run.font.color.rgb = COLOR_DARK_TEXT
                    self._set_font_east_asian(run, FONT_BODY)

            # 正文 — 20pt TencentSans W3
            elif col.get('text'):
                text = col['text']
                if len(text) > 400:
                    text = text[:400] + '...'
                text_lines = text.split('\n')
                for line in text_lines[:12]:
                    if not line.strip():
                        continue
                    p = tf.add_paragraph()
                    p.space_before = Pt(5)
                    p.alignment = PP_ALIGN.LEFT
                    run = p.add_run()
                    run.text = line
                    run.font.name = FONT_BODY
                    run.font.size = Pt(20)
                    run.font.color.rgb = COLOR_DARK_TEXT
                    self._set_font_east_asian(run, FONT_BODY)

            # 高亮框 — 18pt W3 南京主蓝
            if col.get('highlight'):
                hl = col['highlight']
                if len(hl) > 300:
                    hl = hl[:300] + '...'
                hl_lines = hl.split('\n')
                for hl_idx, hl_line in enumerate(hl_lines[:5]):
                    if not hl_line.strip():
                        continue
                    p = tf.add_paragraph()
                    p.space_before = Pt(8)
                    p.alignment = PP_ALIGN.LEFT
                    run = p.add_run()
                    run.text = f"💡 {hl_line}" if hl_idx == 0 else f"   {hl_line}"
                    run.font.name = FONT_BODY
                    run.font.size = Pt(18)
                    run.font.color.rgb = COLOR_BRAND_RED
                    self._set_font_east_asian(run, FONT_BODY)

        return y_start + col_height + Inches(0.3)

    def _render_flow(self, slide, flow_steps: list, y_start) -> int:
        """渲染流程图（步骤文字22pt W7）"""
        if not flow_steps:
            return y_start

        n_steps = len(flow_steps)
        total_width = Inches(17.0)
        step_left_start = Inches(1.2)
        arrow_width = Inches(0.5)
        step_width = int((total_width - arrow_width * (n_steps - 1)) / n_steps)
        step_height = Inches(1.4)

        for idx, step in enumerate(flow_steps):
            left = step_left_start + idx * (step_width + arrow_width)
            
            # 步骤框
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y_start, step_width, step_height)
            shape.fill.solid()
            # 根据风格设置颜色
            step_style = step.get('style', 'default')
            if step_style == 'accent':
                shape.fill.fore_color.rgb = COLOR_WHITE
                shape.line.color.rgb = COLOR_BRAND_RED
            else:
                shape.fill.fore_color.rgb = COLOR_CREAM
                shape.line.color.rgb = COLOR_LINE_GRAY
            shape.adjustments[0] = 0.1

            tf = shape.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            text = step['text']
            if len(text) > 40:
                text = text[:40]
            run.text = text
            run.font.name = FONT_TITLE
            run.font.size = Pt(22)
            run.font.color.rgb = COLOR_BRAND_RED if step_style == 'accent' else COLOR_DARK_TEXT
            run.font.bold = True
            self._set_font_east_asian(run, FONT_TITLE)

            # 箭头（除了最后一步）
            if idx < n_steps - 1:
                arrow_left = left + step_width
                arrow_top = y_start + step_height // 2 - Inches(0.2)
                atf = self._add_textbox(slide, arrow_left, arrow_top, arrow_width, Inches(0.5))
                self._add_text_run(atf, "→", font_name=FONT_BODY, font_size=Pt(28),
                                 color=COLOR_BRAND_RED, bold=True, alignment=PP_ALIGN.CENTER)

        return y_start + step_height + Inches(0.3)

    def _render_highlight_box(self, slide, text: str, y_start) -> int:
        """渲染高亮框（正文20pt W3，左侧红色锚条）"""
        left = Inches(1.0)
        width = Inches(17.5)
        lines = text.count('\n') + 1
        height = Inches(min(max(1.0, lines * 0.45 + 0.4), 3.0))

        # 主背景框
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y_start, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CREAM
        shape.line.color.rgb = COLOR_LINE_GRAY
        shape.adjustments[0] = 0.02

        # 左侧红色锚条
        anchor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y_start, Inches(0.12), height)
        anchor.fill.solid()
        anchor.fill.fore_color.rgb = COLOR_BRAND_RED
        anchor.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.5)
        tf.margin_right = Inches(0.4)
        tf.margin_top = Inches(0.18)
        tf.margin_bottom = Inches(0.18)
        
        if len(text) > 350:
            text = text[:350] + '...'
        # 按行渲染
        text_lines = text.split('\n')
        for i, line in enumerate(text_lines[:8]):
            if not line.strip():
                continue
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(5)
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line
            run.font.name = FONT_BODY
            run.font.size = Pt(20)
            run.font.color.rgb = COLOR_DARK_TEXT
            self._set_font_east_asian(run, FONT_BODY)

        return y_start + height + Inches(0.2)

    def _render_summary_grid(self, slide, summary_cards: list, y_start) -> int:
        """渲染总结网格（方法标签+标题+描述+类比）"""
        return self._render_cards(slide, [
            {
                'icon': '',
                'title': f"[{sc.get('label', '')}] {sc.get('title', '')}",
                'body': sc.get('desc', ''),
                'tag': sc.get('analogy', ''),
                'style': 'default'
            }
            for sc in summary_cards
        ], y_start)

    def _render_policy_grid(self, slide, policy_cards: list, y_start) -> int:
        """渲染策略对比卡片（大emoji + 角色关系 + 标签 + 描述）
        
        对应 HTML 的 .policy-grid > .policy-card 结构：
          - .who: 答题：孩子 / 学习：孩子
          - emoji: 👦→👦 (42px)
          - .label: 同策略 On-Policy
          - desc: 补充说明文字
        """
        n_cards = len(policy_cards)
        if n_cards == 0:
            return y_start

        card_gap = Inches(0.4)
        total_width = Inches(17.0)
        card_left_start = Inches(1.2)
        card_width = int((total_width - card_gap * (n_cards - 1)) / n_cards)
        card_height = Inches(3.8)

        available = Inches(9.5) - y_start
        if card_height > available:
            card_height = available

        for idx, pc in enumerate(policy_cards):
            left = card_left_start + idx * (card_width + card_gap)

            # 卡片背景框（使用不同颜色区分 On-Policy / Off-Policy）
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y_start, card_width, card_height)
            shape.fill.solid()
            label = pc.get('label', '')
            if 'On' in label or '同策略' in label:
                shape.fill.fore_color.rgb = COLOR_WHITE  # 淡蓝
                shape.line.color.rgb = COLOR_PRIMARY_BLUE      # 蓝色边框
            elif 'Off' in label or '异策略' in label:
                shape.fill.fore_color.rgb = COLOR_WHITE  # 暖橙
                shape.line.color.rgb = COLOR_BRAND_RED                   # 南京主蓝强调边框
            else:
                shape.fill.fore_color.rgb = COLOR_CREAM
                shape.line.color.rgb = COLOR_LINE_GRAY
            shape.adjustments[0] = 0.05

            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.4)
            tf.margin_right = Inches(0.4)
            tf.margin_top = Inches(0.3)
            tf.margin_bottom = Inches(0.3)
            tf.vertical_anchor = MSO_ANCHOR.TOP

            # 角色关系（who）— 20pt W3
            who = pc.get('who', '')
            if who:
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = who
                run.font.name = FONT_BODY
                run.font.size = Pt(20)
                run.font.color.rgb = COLOR_MID_GRAY
                self._set_font_east_asian(run, FONT_BODY)

            # Emoji（大字 42pt）
            emoji_text = pc.get('emoji', '')
            if emoji_text:
                p = tf.add_paragraph()
                p.space_before = Pt(10)
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = emoji_text
                run.font.size = Pt(42)

            # 策略标签（label）— 28pt W7 品牌色
            if label:
                p = tf.add_paragraph()
                p.space_before = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = label
                run.font.name = FONT_TITLE
                run.font.size = Pt(28)
                run.font.bold = True
                if 'On' in label or '同策略' in label:
                    run.font.color.rgb = COLOR_PRIMARY_BLUE
                else:
                    run.font.color.rgb = COLOR_BRAND_RED
                self._set_font_east_asian(run, FONT_TITLE)

            # 描述 — 20pt W3
            desc = pc.get('desc', '')
            if desc:
                desc_lines = desc.split('\n') if '\n' in desc else [desc]
                for dl in desc_lines[:4]:
                    if not dl.strip():
                        continue
                    p = tf.add_paragraph()
                    p.space_before = Pt(8)
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = dl
                    run.font.name = FONT_BODY
                    run.font.size = Pt(20)
                    run.font.color.rgb = COLOR_MID_GRAY
                    self._set_font_east_asian(run, FONT_BODY)

        return y_start + card_height + Inches(0.3)

    def _render_prob_chart(self, slide, chart: dict, y_start) -> int:
        """渲染概率条形图（模拟 token 概率分布可视化）"""
        items = chart.get('items', [])
        if not items:
            return y_start

        chart_title = chart.get('title', '概率分布')
        left = Inches(1.0)
        width = Inches(8.5)
        
        # 标题
        height = Inches(0.6)
        tf = self._add_textbox(slide, left, y_start, width, height)
        self._add_text_run(tf, chart_title, font_name=FONT_TITLE, font_size=Pt(20),
                         color=COLOR_DARK_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
        y_cursor = y_start + Inches(0.6)

        # 条形
        bar_left = Inches(1.0)
        bar_width = Inches(8.0)
        bar_height = Inches(0.45)
        label_width = Inches(1.0)
        pct_width = Inches(0.8)

        for item in items[:6]:  # 最多6行
            label = item.get('label', '')
            pct_text = item.get('pct', '0%')
            
            # 标签
            tf = self._add_textbox(slide, bar_left, y_cursor, label_width, bar_height)
            self._add_text_run(tf, label, font_name=FONT_TITLE, font_size=Pt(18),
                             color=COLOR_DARK_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
            
            # 条形背景
            bar_x = bar_left + label_width + Inches(0.1)
            actual_bar_width = bar_width - label_width - pct_width - Inches(0.2)
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                             bar_x, y_cursor + Inches(0.08),
                                             actual_bar_width, Inches(0.28))
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = COLOR_LINE_GRAY
            bg_shape.line.fill.background()
            bg_shape.adjustments[0] = 0.3

            # 填充条
            try:
                pct_val = int(pct_text.replace('%', ''))
            except ValueError:
                pct_val = 10
            fill_width = int(actual_bar_width * pct_val / 100)
            if fill_width > 0:
                fill_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                   bar_x, y_cursor + Inches(0.08),
                                                   fill_width, Inches(0.28))
                fill_shape.fill.solid()
                fill_shape.fill.fore_color.rgb = COLOR_BRAND_RED if pct_val > 30 else COLOR_WARM_APRICOT
                fill_shape.line.fill.background()
                fill_shape.adjustments[0] = 0.3

            # 百分比文字
            pct_x = bar_x + actual_bar_width + Inches(0.1)
            tf = self._add_textbox(slide, pct_x, y_cursor, pct_width, bar_height)
            self._add_text_run(tf, pct_text, font_name=FONT_TITLE, font_size=Pt(16),
                             color=COLOR_BRAND_RED if pct_val > 30 else COLOR_MID_GRAY,
                             bold=True, alignment=PP_ALIGN.RIGHT)

            y_cursor += bar_height

        return y_cursor + Inches(0.2)

    def _render_info_cards(self, slide, info_cards: list, y_start) -> int:
        """渲染信息卡片（采样方法等小型说明卡）"""
        if not info_cards:
            return y_start

        n_cards = min(len(info_cards), 3)
        card_gap = Inches(0.2)
        total_width = Inches(8.5)
        card_left = Inches(10.0)  # 放在右侧
        card_width = total_width
        card_height = Inches(1.5)

        for idx, ic in enumerate(info_cards[:3]):
            top = y_start + idx * (card_height + card_gap)
            if top > Inches(9.0):
                break

            # 卡片框
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         card_left, top, card_width, card_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_CREAM
            shape.line.color.rgb = COLOR_BRAND_RED if ic.get('style') == 'accent' else COLOR_LINE_GRAY
            shape.adjustments[0] = 0.04

            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.25)
            tf.margin_right = Inches(0.25)
            tf.margin_top = Inches(0.15)
            tf.margin_bottom = Inches(0.15)
            tf.vertical_anchor = MSO_ANCHOR.TOP

            # 标题 — 20pt W7
            if ic.get('title'):
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = ic['title']
                run.font.name = FONT_TITLE
                run.font.size = Pt(20)
                run.font.bold = True
                run.font.color.rgb = COLOR_BRAND_RED if ic.get('style') == 'accent' else COLOR_DARK_TEXT
                self._set_font_east_asian(run, FONT_TITLE)

            # 正文 — 18pt W3
            if ic.get('body'):
                body = ic['body']
                if len(body) > 150:
                    body = body[:150] + '...'
                p = tf.add_paragraph()
                p.space_before = Pt(6)
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = body
                run.font.name = FONT_BODY
                run.font.size = Pt(16)
                run.font.color.rgb = COLOR_DARK_TEXT
                self._set_font_east_asian(run, FONT_BODY)

        return y_start + n_cards * (card_height + card_gap)

    def _render_gpt_letters(self, slide, letters: list, y_start) -> int:
        """渲染 GPT 大字展示"""
        if not letters:
            return y_start

        # 大字居中
        total_width = Inches(16.0)
        left = Inches(2.0)
        height = Inches(2.5)
        
        tf = self._add_textbox(slide, left, y_start, total_width, height)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        for idx, letter_data in enumerate(letters):
            run = p.add_run()
            run.text = letter_data.get('letter', '')
            run.font.name = FONT_TITLE
            run.font.size = Pt(120)
            run.font.bold = True
            # 交替颜色
            colors = [COLOR_BRAND_RED, COLOR_WARM_APRICOT, COLOR_DARK_TEXT]
            run.font.color.rgb = colors[idx % len(colors)]
            self._set_font_east_asian(run, FONT_TITLE)
            # 字间距
            if idx < len(letters) - 1:
                spacer = p.add_run()
                spacer.text = '  '
                spacer.font.size = Pt(80)

        return y_start + height + Inches(0.3)

    def _render_inline_flex(self, slide, items: list, y_start) -> int:
        """渲染内联 flex 块（如 GPT 字母解释的三栏说明）"""
        if not items:
            return y_start

        n_items = len(items)
        gap = Inches(0.3)
        total_width = Inches(17.0)
        item_left_start = Inches(1.5)
        item_width = int((total_width - gap * (n_items - 1)) / n_items)
        item_height = Inches(3.0)

        available = Inches(9.5) - y_start
        if item_height > available:
            item_height = available

        for idx, item in enumerate(items):
            left = item_left_start + idx * (item_width + gap)

            # 卡片框
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         left, y_start, item_width, item_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_CREAM
            shape.line.color.rgb = COLOR_LINE_GRAY
            shape.adjustments[0] = 0.04

            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.25)
            tf.margin_right = Inches(0.25)
            tf.margin_top = Inches(0.2)
            tf.margin_bottom = Inches(0.2)
            tf.vertical_anchor = MSO_ANCHOR.TOP

            # 图标
            if item.get('icon'):
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = item['icon']
                run.font.size = Pt(36)

            # 标题 — 28pt W7
            if item.get('title'):
                p = tf.add_paragraph()
                p.space_before = Pt(10)
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = item['title']
                run.font.name = FONT_TITLE
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = COLOR_BRAND_RED
                self._set_font_east_asian(run, FONT_TITLE)

            # 描述 — 20pt W3
            if item.get('desc'):
                desc = item['desc']
                for line in desc.split('\n')[:4]:
                    p = tf.add_paragraph()
                    p.space_before = Pt(6)
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = line
                    run.font.name = FONT_BODY
                    run.font.size = Pt(18)
                    run.font.color.rgb = COLOR_DARK_TEXT
                    self._set_font_east_asian(run, FONT_BODY)

        return y_start + item_height + Inches(0.3)

    def generate(self, blocks: List[ContentBlock], title: str = '', subtitle: str = '', author: str = '') -> None:
        """从内容块列表生成完整 PPT（文章型 HTML）"""
        # 1. 生成封面页
        self._add_cover_slide(title, subtitle, author)
        
        # 2. 遍历内容块，智能分页生成
        current_content: List[ContentBlock] = []
        current_title = ''
        
        i = 0
        while i < len(blocks):
            block = blocks[i]
            
            # H1 → 生成章节扇页
            if block.type == 'h1':
                # 先输出已累积的内容
                if current_content:
                    self._flush_content_slides(current_title, current_content)
                    current_content = []
                # 章节页
                self._add_section_slide(block.text)
                current_title = ''
                i += 1
                continue
            
            # H2 → 新内容页标题
            if block.type == 'h2':
                # 先输出已累积的内容
                if current_content:
                    self._flush_content_slides(current_title, current_content)
                    current_content = []
                current_title = block.text
                i += 1
                continue
            
            # 表格 → 独立页面
            if block.type == 'table':
                if current_content:
                    self._flush_content_slides(current_title, current_content)
                    current_content = []
                self._add_table_slides(current_title, block)
                current_title = ''
                i += 1
                continue
            
            # 代码块 → 独立页面
            if block.type == 'code':
                if current_content:
                    self._flush_content_slides(current_title, current_content)
                    current_content = []
                self._add_code_slide(current_title, block)
                current_title = ''
                i += 1
                continue
            
            # 图片 → 独立页面
            if block.type == 'image':
                if current_content:
                    self._flush_content_slides(current_title, current_content)
                    current_content = []
                self._add_image_slide(current_title, block)
                current_title = ''
                i += 1
                continue
            
            # 其他内容累积
            current_content.append(block)
            i += 1
        
        # 输出剩余内容
        if current_content:
            self._flush_content_slides(current_title, current_content)
        
        # 3. 生成结尾页
        self._add_end_slide()
    
    def save(self, output_path: str):
        """保存 PPT 文件"""
        self.prs.save(output_path)
        print(f"[OK] 已生成 PPT: {output_path}")
        print(f"     共 {self.slide_count} 页")
    
    # ---------- 页面类型生成 ----------
    
    def _add_cover_slide(self, title: str, subtitle: str = '', author: str = ''):
        """封面页"""
        slide = self._new_slide('cover')
        
        # 主标题
        if title:
            left = Inches(2.0)
            top = Inches(3.5)
            width = Inches(16.0)
            height = Inches(2.5)
            tf = self._add_textbox(slide, left, top, width, height)
            self._add_text_run(tf, title, font_name=FONT_TITLE, font_size=Pt(60),
                             color=COLOR_DARK_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
        
        # 副标题
        if subtitle:
            left = Inches(2.0)
            top = Inches(6.2)
            width = Inches(14.0)
            height = Inches(1.5)
            tf = self._add_textbox(slide, left, top, width, height)
            self._add_text_run(tf, subtitle, font_name=FONT_BODY, font_size=Pt(28),
                             color=COLOR_MID_GRAY, alignment=PP_ALIGN.LEFT)
        
        # 作者
        if author:
            left = Inches(2.0)
            top = Inches(8.0)
            width = Inches(10.0)
            height = Inches(1.0)
            tf = self._add_textbox(slide, left, top, width, height)
            self._add_text_run(tf, author, font_name=FONT_BODY, font_size=Pt(22),
                             color=COLOR_MID_GRAY, alignment=PP_ALIGN.LEFT)
    
    def _add_section_slide(self, title: str):
        """章节扇页"""
        slide = self._new_slide('section')
        
        # 章节标题（居中偏左）
        left = Inches(2.0)
        top = Inches(4.0)
        width = Inches(16.0)
        height = Inches(3.0)
        tf = self._add_textbox(slide, left, top, width, height)
        self._add_text_run(tf, title, font_name=FONT_TITLE, font_size=Pt(52),
                         color=COLOR_DARK_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
    
    def _add_end_slide(self):
        """结尾页"""
        slide = self._new_slide('end')
        
        left = Inches(5.0)
        top = Inches(4.0)
        width = Inches(10.0)
        height = Inches(3.0)
        tf = self._add_textbox(slide, left, top, width, height)
        self._add_text_run(tf, "THANKS", font_name=FONT_TITLE, font_size=Pt(72),
                         color=COLOR_BRAND_RED, bold=True, alignment=PP_ALIGN.CENTER)
        
        # 副标题
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(20)
        run = p.add_run()
        run.text = "感谢聆听"
        run.font.name = FONT_BODY
        run.font.size = Pt(28)
        run.font.color.rgb = COLOR_MID_GRAY
        self._set_font_east_asian(run, FONT_BODY)
    
    def _flush_content_slides(self, title: str, blocks: List[ContentBlock]):
        """将累积的内容块分页输出为内容页"""
        # 估算内容量，决定分页
        pages = self._paginate_content(blocks)
        
        for page_blocks in pages:
            slide = self._new_slide('content')
            y_cursor = Inches(1.5)
            
            # 页面标题 — 36pt TencentSans W7
            if title:
                left = Inches(1.0)
                width = Inches(17.0)
                height = Inches(1.0)
                tf = self._add_textbox(slide, left, y_cursor, width, height)
                self._add_text_run(tf, title, font_name=FONT_TITLE, font_size=Pt(36),
                                 color=COLOR_DARK_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
                y_cursor += Inches(1.2)
            
            # 内容渲染
            for block in page_blocks:
                y_cursor = self._render_content_block(slide, block, y_cursor)
                if y_cursor > Inches(9.5):
                    break  # 超出安全区域，停止
    
    def _add_table_slides(self, title: str, block: ContentBlock):
        """表格页（可能分多页）"""
        headers = block.headers
        rows = block.rows
        
        # 分页
        row_chunks = [rows[i:i+MAX_TABLE_ROWS_PER_SLIDE] for i in range(0, max(len(rows), 1), MAX_TABLE_ROWS_PER_SLIDE)]
        
        for chunk_idx, chunk in enumerate(row_chunks):
            slide = self._new_slide('content')
            y_cursor = Inches(1.5)
            
            # 标题
            page_title = title if title else "数据表格"
            if len(row_chunks) > 1:
                page_title += f" ({chunk_idx + 1}/{len(row_chunks)})"
            
            left = Inches(1.0)
            width = Inches(17.0)
            height = Inches(1.0)
            tf = self._add_textbox(slide, left, y_cursor, width, height)
            self._add_text_run(tf, page_title, font_name=FONT_TITLE, font_size=Pt(32),
                             color=COLOR_DARK_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
            y_cursor += Inches(1.2)
            
            # 生成表格
            n_rows = len(chunk) + (1 if headers else 0)
            n_cols = len(headers) if headers else (len(chunk[0]) if chunk else 1)
            
            table_left = Inches(1.0)
            table_width = Inches(18.0)
            row_height = Inches(0.65)
            table_height = Emu(int(n_rows * row_height))
            
            table_shape = slide.shapes.add_table(n_rows, n_cols, table_left, y_cursor, table_width, table_height)
            table = table_shape.table
            
            # 设置列宽
            col_width = int(table_width / n_cols)
            for col_idx in range(n_cols):
                table.columns[col_idx].width = col_width
            
            # 填充表头
            row_idx = 0
            if headers:
                for col_idx, header_text in enumerate(headers[:n_cols]):
                    cell = table.cell(0, col_idx)
                    cell.text = header_text
                    self._style_table_cell(cell, is_header=True)
                row_idx = 1
            
            # 填充数据
            for data_row in chunk:
                for col_idx, cell_text in enumerate(data_row[:n_cols]):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = cell_text
                    self._style_table_cell(cell, is_header=False)
                row_idx += 1
    
    def _add_code_slide(self, title: str, block: ContentBlock):
        """代码块页"""
        slide = self._new_slide('content')
        y_cursor = Inches(1.5)
        
        # 标题
        page_title = title if title else (f"代码示例" + (f" ({block.language})" if block.language else ""))
        left = Inches(1.0)
        width = Inches(17.0)
        height = Inches(1.0)
        tf = self._add_textbox(slide, left, y_cursor, width, height)
        self._add_text_run(tf, page_title, font_name=FONT_TITLE, font_size=Pt(32),
                         color=COLOR_DARK_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
        y_cursor += Inches(1.2)
        
        # 代码框（深色背景）
        code_left = Inches(1.0)
        code_width = Inches(18.0)
        code_height = Inches(7.5)
        
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, code_left, y_cursor, code_width, code_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BLACK
        shape.line.fill.background()
        # 圆角
        shape.adjustments[0] = 0.02
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.4)
        tf.margin_right = Inches(0.4)
        tf.margin_top = Inches(0.3)
        tf.margin_bottom = Inches(0.3)
        
        # 截断过长代码
        code_lines = block.text.split('\n')
        if len(code_lines) > 25:
            code_lines = code_lines[:25] + ['...']
        code_text = '\n'.join(code_lines)
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = code_text
        run.font.name = FONT_CODE
        run.font.size = Pt(14)
        run.font.color.rgb = COLOR_WHITE
        # 设置等宽字体
        rPr = run._r.get_or_add_rPr()
        latin = rPr.find(qn('a:latin'))
        if latin is None:
            latin = etree.SubElement(rPr, qn('a:latin'))
        latin.set('typeface', FONT_CODE)
    
    def _add_image_slide(self, title: str, block: ContentBlock):
        """图片页"""
        slide = self._new_slide('content')
        y_cursor = Inches(1.5)
        
        # 标题
        if title or block.alt:
            page_title = title or block.alt
            left = Inches(1.0)
            width = Inches(17.0)
            height = Inches(1.0)
            tf = self._add_textbox(slide, left, y_cursor, width, height)
            self._add_text_run(tf, page_title, font_name=FONT_TITLE, font_size=Pt(32),
                             color=COLOR_DARK_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
            y_cursor += Inches(1.2)
        
        # 尝试加载图片
        img_path = self._resolve_image(block.src)
        if img_path and os.path.exists(img_path):
            # 居中放置图片
            img_left = Inches(3.0)
            img_top = y_cursor
            img_width = Inches(14.0)
            img_height = Inches(7.5)
            try:
                slide.shapes.add_picture(img_path, img_left, img_top, width=img_width)
            except Exception:
                # 图片加载失败，添加占位文字
                self._add_image_placeholder(slide, y_cursor, block.src)
        else:
            self._add_image_placeholder(slide, y_cursor, block.src)
    
    def _add_image_placeholder(self, slide, y_cursor, src: str):
        """图片占位符"""
        left = Inches(4.0)
        width = Inches(12.0)
        height = Inches(6.0)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y_cursor, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_LINE_GRAY
        shape.line.fill.background()
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"[图片: {src}]"
        run.font.name = FONT_BODY
        run.font.size = Pt(18)
        run.font.color.rgb = COLOR_MID_GRAY
    
    # ---------- 内容渲染 ----------
    
    def _render_content_block(self, slide, block: ContentBlock, y_cursor) -> int:
        """渲染单个内容块到幻灯片，返回新的 y_cursor
        
        字体规范：
          - h3: 28pt W7, h4-h6: 24pt W7
          - paragraph: 20pt W3
          - list: 20pt W3
          - blockquote: 20pt W3
        """
        
        if block.type in ('h3', 'h4', 'h5', 'h6'):
            # 副标题 — 24-28pt TencentSans W7
            left = Inches(1.0)
            width = Inches(17.0)
            height = Inches(0.8)
            tf = self._add_textbox(slide, left, y_cursor, width, height)
            size = Pt(28) if block.type == 'h3' else Pt(24)
            self._add_text_run(tf, block.text, font_name=FONT_TITLE, font_size=size,
                             color=COLOR_DARK_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
            return y_cursor + Inches(1.0)
        
        elif block.type == 'paragraph':
            # 正文段落 — 20pt TencentSans W3
            left = Inches(1.0)
            width = Inches(17.0)
            lines = math.ceil(len(block.text) / 60)
            height = Inches(max(0.6, lines * 0.45))
            tf = self._add_textbox(slide, left, y_cursor, width, height)
            self._add_text_run(tf, block.text, font_name=FONT_BODY, font_size=Pt(20),
                             color=COLOR_DARK_TEXT, alignment=PP_ALIGN.LEFT)
            return y_cursor + height + Inches(0.2)
        
        elif block.type == 'list':
            # 列表 — 20pt TencentSans W3
            left = Inches(1.2)
            width = Inches(16.8)
            
            for idx, item in enumerate(block.items):
                if y_cursor > Inches(9.5):
                    break
                height = Inches(0.55)
                tf = self._add_textbox(slide, left, y_cursor, width, height)
                if item.startswith('  • '):
                    prefix = '    • '
                    text = item[4:]
                elif block.ordered:
                    prefix = f'{idx + 1}. '
                    text = item
                else:
                    prefix = '• '
                    text = item
                
                self._add_text_run(tf, prefix + text, font_name=FONT_BODY, font_size=Pt(20),
                                 color=COLOR_DARK_TEXT, alignment=PP_ALIGN.LEFT)
                y_cursor += Inches(0.55)
            
            return y_cursor + Inches(0.2)
        
        elif block.type == 'blockquote':
            # 引用框 — 20pt W3 + 左侧红色锚条
            left = Inches(1.0)
            width = Inches(17.0)
            lines = math.ceil(len(block.text) / 55)
            height = Inches(max(1.0, lines * 0.45 + 0.4))
            
            # 背景框
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y_cursor, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_CREAM
            shape.line.fill.background()
            shape.adjustments[0] = 0.02
            
            # 左侧红色锚条
            anchor_left = left
            anchor_width = Inches(0.12)
            anchor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, anchor_left, y_cursor, anchor_width, height)
            anchor.fill.solid()
            anchor.fill.fore_color.rgb = COLOR_BRAND_RED
            anchor.line.fill.background()
            
            # 文字
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.5)
            tf.margin_top = Inches(0.2)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = block.text
            run.font.name = FONT_BODY
            run.font.size = Pt(20)
            run.font.color.rgb = COLOR_MID_GRAY
            self._set_font_east_asian(run, FONT_BODY)
            
            return y_cursor + height + Inches(0.3)
        
        return y_cursor
    
    # ---------- 分页逻辑 ----------
    
    def _paginate_content(self, blocks: List[ContentBlock]) -> List[List[ContentBlock]]:
        """智能分页：根据内容量将块分配到多页"""
        pages = []
        current_page = []
        current_lines = 0
        
        for block in blocks:
            est_lines = self._estimate_lines(block)
            
            # 如果当前页已满，开新页
            if current_lines + est_lines > MAX_LINES_PER_SLIDE and current_page:
                pages.append(current_page)
                current_page = []
                current_lines = 0
            
            current_page.append(block)
            current_lines += est_lines
        
        if current_page:
            pages.append(current_page)
        
        return pages if pages else [[]]
    
    def _estimate_lines(self, block: ContentBlock) -> int:
        """估算内容块占用的行数"""
        if block.type in ('h3', 'h4', 'h5', 'h6'):
            return 2
        elif block.type == 'paragraph':
            return max(1, math.ceil(len(block.text) / 60)) + 1
        elif block.type == 'list':
            return len(block.items) + 1
        elif block.type == 'blockquote':
            return max(2, math.ceil(len(block.text) / 55) + 1)
        return 1
    
    # ---------- 底层辅助 ----------
    
    def _new_slide(self, page_type: str):
        """创建新幻灯片并设置背景+Logo"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self.slide_count += 1
        
        # 设置背景
        bg_path = {'cover': BG_COVER, 'section': BG_SECTION, 'end': BG_END}.get(page_type, BG_CONTENT)
        self._set_background(slide, str(bg_path))
        
        # 南京模板默认不强制插入旧横版 Logo
        self._add_logo(slide, page_type)
        
        return slide
    
    def _set_background(self, slide, image_path: str):
        """设置幻灯片背景图（铺满整页，z-order 最底层）"""
        if not os.path.exists(image_path):
            return
        
        pic = slide.shapes.add_picture(image_path, 0, 0, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU)
        # 移到最底层
        sp_tree = slide.shapes._spTree
        pic_elem = pic._element
        sp_tree.remove(pic_elem)
        sp_tree.insert(2, pic_elem)  # 在 spPr 后面
    
    def _add_logo(self, slide, page_type: str):
        """南京模板默认不强制插入横版 Logo；保留钩子兼容未来官方 Logo。"""
        return
    
    def _add_textbox(self, slide, left, top, width, height):
        """添加文本框，返回 text_frame"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        return tf
    
    def _add_text_run(self, text_frame, text: str, font_name: str = FONT_BODY,
                     font_size=Pt(20), color=COLOR_DARK_TEXT, bold: bool = False,
                     alignment=PP_ALIGN.LEFT):
        """向 text_frame 添加格式化文本"""
        p = text_frame.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.bold = bold
        # 设置中文字体
        self._set_font_east_asian(run, font_name)
        return run
    
    def _set_font_east_asian(self, run, font_name: str):
        """设置 East Asian 字体（中文字体必须通过 a:ea 设置）"""
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', font_name)
        # 同时设置 latin
        latin = rPr.find(qn('a:latin'))
        if latin is None:
            latin = etree.SubElement(rPr, qn('a:latin'))
        latin.set('typeface', font_name)
    
    def _style_table_cell(self, cell, is_header: bool = False):
        """设置表格单元格样式"""
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = FONT_TITLE if is_header else FONT_BODY
                run.font.size = Pt(16) if is_header else Pt(14)
                run.font.color.rgb = COLOR_WHITE if is_header else COLOR_DARK_TEXT
                run.font.bold = is_header
                self._set_font_east_asian(run, FONT_TITLE if is_header else FONT_BODY)
        
        # 单元格填充
        tcPr = cell._tc.get_or_add_tcPr()
        # 设置填充
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        if is_header:
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', '3272DC')  # 南京主蓝表头
        else:
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', 'FFFFFF')  # 白色数据行
    
    def _resolve_image(self, src: str) -> Optional[str]:
        """解析图片路径/URL"""
        if not src:
            return None
        
        # 本地路径
        if os.path.exists(src):
            return src
        
        # 相对路径
        # URL 下载
        if src.startswith(('http://', 'https://')):
            try:
                tmp_dir = tempfile.mkdtemp()
                ext = os.path.splitext(src)[1][:5] or '.png'
                tmp_path = os.path.join(tmp_dir, f'img{ext}')
                urllib.request.urlretrieve(src, tmp_path)
                return tmp_path
            except Exception:
                return None
        
        return None


# ---------- 主入口 ----------

def main():
    parser = argparse.ArgumentParser(description='HTML → 腾讯云架构师南京城市沙龙模板 PPT')
    parser.add_argument('--input', '-i', required=True, help='输入 HTML 文件路径')
    parser.add_argument('--output', '-o', required=True, help='输出 PPTX 文件路径')
    parser.add_argument('--title', '-t', default='', help='PPT 主标题（默认从 HTML title/h1 提取）')
    parser.add_argument('--subtitle', '-s', default='', help='PPT 副标题')
    parser.add_argument('--author', '-a', default='', help='作者/演讲者')
    parser.add_argument('--encoding', default='utf-8', help='HTML 文件编码（默认 utf-8）')
    
    args = parser.parse_args()
    
    # 读取 HTML
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"[ERROR] 输入文件不存在: {args.input}", file=sys.stderr)
        sys.exit(1)
    
    html_content = input_path.read_text(encoding=args.encoding)
    
    # 检测 HTML 类型
    print(f"[INFO] 解析 HTML: {args.input}")
    
    if is_slide_deck_html(html_content):
        # 幻灯片型 HTML
        print(f"[INFO] 检测到幻灯片型 HTML（slide deck），按页面结构解析")
        slide_contents = parse_html_slides(html_content)
        print(f"[INFO] 解析得到 {len(slide_contents)} 个幻灯片页面")
        
        gen = PptxGenerator()
        gen.generate_from_slides(slide_contents)
    else:
        # 文章型 HTML
        print(f"[INFO] 检测到文章型 HTML，按内容结构解析")
        blocks = parse_html(html_content)
        print(f"[INFO] 解析得到 {len(blocks)} 个内容块")
        
        # 提取标题（如果用户没有指定）
        title = args.title
        if not title:
            soup = BeautifulSoup(html_content, 'lxml')
            title_tag = soup.find('title')
            if title_tag:
                title = title_tag.get_text(strip=True)
            elif blocks and blocks[0].type == 'h1':
                title = blocks[0].text
                blocks = blocks[1:]
        
        gen = PptxGenerator()
        gen.generate(blocks, title=title, subtitle=args.subtitle, author=args.author)
    
    # 保存
    os.makedirs(os.path.dirname(os.path.abspath(args.output)), exist_ok=True)
    gen.save(args.output)
    
    # 输出统计
    print(f"\n[报告]")
    print(f"  输入: {args.input}")
    print(f"  输出: {args.output}")
    print(f"  页数: {gen.slide_count}")
    if is_slide_deck_html(html_content):
        print(f"  模式: 幻灯片型 HTML → 结构化迁移")
        print(f"  解析幻灯片: {len(slide_contents)} 页")
    else:
        print(f"  模式: 文章型 HTML → 内容分页")
        print(f"  标题: {title}")
        print(f"  内容块: {len(blocks)}")
    print(f"  字体规范: 标题36-44pt {FONT_TITLE} / 卡片标题24pt / 正文20pt {FONT_BODY}")
    print(f"  品牌合规: 配色=品牌安全色, 背景=模板背景, Logo=南京模板不强制插入")


if __name__ == '__main__':
    main()
